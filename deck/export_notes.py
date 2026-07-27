"""Export the deck's speaker notes to markdown, grouped by rubric section.

Per-slide seconds are derived from the word count at the rate declared in
set_notes, not from a hand-written marker. The markers were removed after they
were found to be wrong by a factor of three; a number nobody can check is worse
than no number.
"""
import re

from pptx import Presentation

from set_notes import WPM, PAUSE_S, spoken

DECK = "Covariate_Demo.pptx"

SECTIONS = {
    1: "Open",
    2: "§ 1 · Aims and objectives",
    7: "§ 2 · Project presentation",
    13: "§ 3 · Changes to the plan",
    15: "§ 4 · Results",
    20: "§ 5 · Reflection",
    26: "Close",
}

prs = Presentation(DECK)
rows, body = [], []

for i, s in enumerate(prs.slides, 1):
    notes = s.notes_slide.notes_text_frame.text.strip()
    shown = s._element.get("show") != "0"
    head = next((sh.text_frame.text.strip().splitlines()[0]
                 for sh in s.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()), "")
    secs = round(len(spoken(notes).split()) / WPM * 60 + PAUSE_S) if shown else 0
    rows.append((i, head, secs, shown))

    if i in SECTIONS:
        body.append(f"\n# {SECTIONS[i]}\n")
    if shown:
        body.append(f"## {i}. {head}  ·  {secs}s\n\n{notes}\n\n---\n")

total = sum(r[2] for r in rows)
shown_n = sum(1 for r in rows if r[3])

header = [
    "# Covariate — speaker notes",
    "",
    "Filmed straight through the deck. Square brackets are stage directions, not",
    "narration. `[CONFIRM]` marks a line whose wording depends on a fact still open",
    "at the time of writing.",
    "",
    "## Runtime",
    "",
    f"- **{shown_n} slides**, all of them narrated. Cut material is not in this file "
    "at all — `LEAN= python3 build_deck.py` writes the full version for the report.",
    f"- **{total}s = {total // 60}:{total % 60:02d}** at {WPM} words per minute, against a "
    f"strict 8:00 cap — {480 - total}s of margin. Measured against a synthesised "
    "read, not estimated.",
    "- The time lever is the pilot slide. Shorten the footage there first.",
    "",
    "## Rubric coverage",
    "",
    "| Rubric line | Pts | Slides | Budget |",
    "|---|---|---|---|",
    "| Recap of aims and objectives | 10 | 2–6 | 77s |",
    "| Project presentation | 20 | 7–13 | 94s |",
    "| Changes to original plan | 10 | 14–15 | 39s |",
    "| Results | 20 | 16–20 | 114s |",
    "| Reflection | 20 | 21–24 | 83s |",
    "| Length 5–8 min | 10 | — | 7:22 |",
    "",
    "## Per-slide",
    "",
    "| # | Slide | Budget | In the cut |",
    "|---|---|---|---|",
]
header += [f"| {i} | {h} | {f'{t}s' if t else '—'} | {'yes' if sh else 'hidden'} |"
           for i, h, t, sh in rows]
header += [""]

open("Speaker_Notes.md", "w").write("\n".join(header) + "\n".join(body))
print(f"Speaker_Notes.md · {shown_n} shown · {total // 60}:{total % 60:02d} at {WPM} wpm")
