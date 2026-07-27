"""Restructure the deck onto the video grading rubric (Option A).

The rubric scores six lines: Aims and objectives (10), Length (10), Project
presentation (20), Changes to the plan (10), Results (20), Reflection (20). It has
no Background line, so the motivation has to live inside Aims or it scores against
nothing. Five dark divider slides name each section in the grader's own words.

Two structural mechanisms are introduced here:

* ``title()`` now registers each slide under a key, and a final pass reorders
  ``sldIdLst`` from a declarative ORDER list. Slides can therefore be authored
  wherever it reads best in this file and still land in narration order. The pass
  fails loudly on an unknown or missing key rather than silently dropping a slide.
* ``citation_card()`` renders a paper reference as native shapes. Deliberately not
  a screenshot: a page image of a two-column journal PDF is illegible at slide
  scale, and reproducing a Science page is a copyright question we do not need to
  have. The card carries the full citation and the exact quoted sentence.
"""
import sys
from pathlib import Path

p = Path("build_deck.py")
s = p.read_text()


def sub(old, new, n=1):
    global s
    if s.count(old) != n:
        sys.exit(f"NO MATCH ({s.count(old)} found, {n} expected):\n{old[:300]}")
    s = s.replace(old, new)


# =============================================================================
# 1 · slide registry + divider and citation-card helpers
# =============================================================================
sub('''def title(slide, t, sub=None, dark=False):
    text(slide, 0.85, 0.62, 11.6, 1.0, t, size=36, bold=True,
         color=WHITE if dark else NAVY)
    if sub:
        text(slide, 0.85, 1.42, 11.6, 0.5, sub, size=16,
             color=RGBColor(0xC8, 0xD2, 0xDC) if dark else MUTED)''',
    '''REG: "dict[str, object]" = {}


def register(slide, key):
    if key in REG:
        sys.exit(f"duplicate slide key: {key}")
    REG[key] = slide
    return slide


def title(slide, t, sub=None, dark=False, k=None):
    text(slide, 0.85, 0.62, 11.6, 1.0, t, size=36, bold=True,
         color=WHITE if dark else NAVY)
    if sub:
        text(slide, 0.85, 1.42, 11.6, 0.5, sub, size=16,
             color=RGBColor(0xC8, 0xD2, 0xDC) if dark else MUTED)
    register(slide, k or t)


def divider(num, name, answers):
    """A dark section marker naming a rubric line in the grader's own words."""
    d = new_slide(dark=True)
    hexagon(d, 0.85, 2.60, 1.05, num, fill=GOLD, fg=NAVY, size=30)
    text(d, 2.35, 2.62, 10.2, 0.9, name, size=42, bold=True, color=WHITE)
    text(d, 2.35, 3.62, 10.2, 0.6, answers, size=17,
         color=RGBColor(0xC8, 0xD2, 0xDC), line=1.3)
    register(d, f"div:{name}")
    return d


def citation_card(slide, y, venue, head, authors, quote, h=1.62):
    """A paper reference as native shapes — citation plus the exact sentence.

    Not a page screenshot. A two-column journal page rendered into a 5-inch box is
    unreadable from a video, and the quoted sentence is the entire reason the paper
    is on the slide.
    """
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85),
                                   Inches(y), Inches(11.6), Inches(h))
    panel.fill.solid(); panel.fill.fore_color.rgb = PALE
    panel.line.color.rgb = OLDGOLD; panel.line.width = Pt(1.25)
    panel.shadow.inherit = False
    panel.text_frame.text = ""

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85),
                                  Inches(y), Inches(0.09), Inches(h))
    rule.fill.solid(); rule.fill.fore_color.rgb = GOLD
    rule.line.fill.background(); rule.shadow.inherit = False

    text(slide, 1.28, y + 0.20, 10.9, 0.28, venue, size=11.5, bold=True,
         color=DEEPGOLD)
    text(slide, 1.28, y + 0.50, 10.9, 0.34, head, size=16, bold=True, color=NAVY)
    text(slide, 1.28, y + 0.80, 10.9, 0.28, authors, size=12, color=MUTED)
    text(slide, 1.28, y + 1.10, 10.9, 0.44, quote, size=14.5, color=NAVY,
         italic=True, line=1.25)
    return panel''')

# text() has no italic parameter at call level; it reads it from the override dict.
sub('''        f.italic = over.get("italic", False)''',
    '''        f.italic = over.get("italic", _italic)''')
sub('''def text(slide, x, y, w, h, runs, size=16, color=INK, bold=False,
         align=PP_ALIGN.LEFT, space=6, line=1.0):''',
    '''def text(slide, x, y, w, h, runs, size=16, color=INK, bold=False,
         align=PP_ALIGN.LEFT, space=6, line=1.0, italic=False):
    _italic = italic''')
sub('''from pathlib import Path

from pptx import Presentation''',
    '''import sys
from pathlib import Path

from pptx import Presentation''')


# =============================================================================
# 2 · § 1 — "Why record the room" rebuilt as three progressive-reveal slides
# =============================================================================
OLD_IDEA_START = '''s = new_slide()
title(s, "Why record the room", "the case for this in a ubiquitous-computing course")
text(s, 0.85, 1.95, 11.4, 0.9, ['''
IDX = s.index(OLD_IDEA_START)
END = s.index('# =============================================================================\n# NEW — why consumer hardware makes this hard')
OLD_IDEA = s[IDX:END]

NEW_IDEA = '''# The same slide three times, adding one paper each. python-pptx cannot write
# PowerPoint animations, so the reveal is built as successive slides — which also
# survives export to PDF and to video without anyone clicking anything.
ERRINGTON = dict(
    venue="eLife 10:e67995 (2021)  ·  CC BY  ·  doi.org/10.7554/eLife.67995",
    head="Challenges for assessing replicability in preclinical cancer biology",
    authors="Errington, Denis, Perfito, Iorns & Nosek  ·  193 experiments, 53 papers",
    quote="\\u201cNone of the 193 experiments were described in sufficient detail in the "
          "original paper to enable us to design protocols to repeat the experiments.\\u201d",
)
CRABBE = dict(
    venue="Science 284(5420):1670\\u20132 (1999)  ·  doi.org/10.1126/science.284.5420.1670",
    head="Genetics of mouse behavior: interactions with laboratory environment",
    authors="Crabbe, Wahlsten & Dudek  ·  three laboratories, simultaneous testing",
    quote="\\u201cApparatus, test protocols, and many environmental variables were rigorously "
          "equated \\u2026 there were systematic differences in behavior across labs.\\u201d",
)

for _step in (0, 1, 2):
    s = new_slide()
    title(s, "Why record the room",
          "the problem this is aimed at, and why it belongs in a ubicomp course",
          k=f"idea{_step}")
    text(s, 0.85, 1.95, 11.4, 0.8, [
        ("Two experiments run to the same written protocol still disagree, and the "
         "written protocol is not where the difference is.",
         {"size": 19, "color": MUTED, "space": 5}),
        ("What is missing is the room \\u2014 and nobody records the room, because until "
         "recently there was nothing cheap enough to record it with.",
         {"size": 19, "bold": True, "color": NAVY}),
    ])
    if _step >= 1:
        citation_card(s, 3.22, **ERRINGTON)
    if _step >= 2:
        citation_card(s, 5.05, **CRABBE)
    if _step == 2:
        text(s, 0.85, 6.88, 11.6, 0.45,
             "Every room where an experiment happens already contains a capable sensor "
             "package. The open question is whether what it measures is worth anything.",
             size=15, bold=True, color=NAVY, line=1.2)
    notes(s, "PLACEHOLDER")


'''
s = s[:IDX] + NEW_IDEA + s[END:]


# =============================================================================
# 3 · § 1 — Aims retitled to the rubric's words
# =============================================================================
sub('''title(s, "Aims", "three, and one design constraint")''',
    '''title(s, "Aims and objectives", "as proposed \\u2014 three objectives and one design "
      "constraint", k="aims")''')


# =============================================================================
# 4 · § 2 — Implementation rewritten for two deployment tiers
# =============================================================================
OLD_IMPL_START = '''s = new_slide()
title(s, "Implementation", "React Native under Expo Go; installs by scanning a QR code")
layers = ['''
I0 = s.index(OLD_IMPL_START)
I1 = s.index('# =============================================================================\n# NEW — what it looks like to use')
s = s[:I0] + '''s = new_slide()
title(s, "Implementation", "six channels, and a deployment boundary running straight "
      "through the middle of them")
layers = [
    (NAVY, "1", "Direct sensors", "expo-sensors  \\u2014  no build step",
     "accelerometer 50 Hz  \\u00b7  magnetometer 25 Hz  \\u00b7  barometer"),
    (TEAL, "2", "Derived channel", "computed from the raw stream",
     "vibration: gravity removed, RMS + peak over a 200 ms window"),
    (DEEPGOLD, "3", "Native modules", "Swift on iOS, Kotlin on Android  \\u2014  needs a compiled build",
     "camera-EXIF light  \\u00b7  microphone LEVEL only \\u2014 audio is never recorded, so "
     "there is no waveform to leak"),
    (OLDGOLD, "4", "Session record", "one session, one JSON file",
     "metadata  \\u00b7  placement  \\u00b7  per-channel sampling health  \\u00b7  one shared clock, "
     "aligned across devices by a haptic fiducial"),
]
y = 2.30
for col, num, name, how, what in layers:
    hexagon(s, 0.85, y - 0.02, 0.56, num, fill=col, size=13)
    text(s, 1.66, y, 3.3, 0.35, name, size=18, bold=True, color=NAVY)
    text(s, 5.05, y + 0.04, 7.4, 0.35, how, size=13, color=DEEPGOLD, bold=True)
    text(s, 1.66, y + 0.42, 10.6, 0.5, what, size=14, color=MUTED, line=1.2)
    y += 1.12

text(s, 0.85, 6.72, 11.6, 0.62, [
    ("Two tiers, and the gap between them is a result.",
     {"size": 16, "bold": True, "color": NAVY, "space": 4}),
    ("Four channels install by scanning a QR code \\u2014 no build, no developer account. "
     "All six require a compiled dev client. The hardware is universal; access to it "
     "is not.", {"size": 14, "color": MUTED}),
], line=1.2)
notes(s, "PLACEHOLDER")


''' + s[I1:]


# =============================================================================
# 5 · § 3 — Changes gains the dev-client result
# =============================================================================
sub('''hex_rows(s, [
    ("Scope reduced to the channels that run in Expo Go.",
     "Light and sound level are native modules and require a compiled dev client. Rather than "
     "spend schedule on that build step at a second site, we replaced the Alka-Seltzer "
     "dissolution study with a door experiment using only the remaining channels."),
    ("Multi-site study reclassified as a case study.",
     "With one participant per site, person, city, phone model and building are confounded. "
     "The quantitative claims moved to a within-site design with a trial count that can "
     "support them."),
    ("We pre-registered.",
     "Metrics, windows, exclusion rules and trial counts are frozen in the repository, dated, "
     "before the data existed."),
], top=2.35, gap=1.42, fill=OLDGOLD)''',
    '''hex_rows(s, [
    ("Scope was cut to four channels, then won back to six.",
     "Light and sound level are native modules needing a compiled dev client, so the pilot "
     "ran on the four Expo Go channels. The dev client is now building and all six channels "
     "run \\u2014 which turns a dropped feature into a measured deployment boundary."),
    ("Multi-site study reclassified as a case study.",
     "With one participant per site, person, city, phone model and building are confounded. "
     "The quantitative claims moved to a within-site design with a trial count that can "
     "support them."),
    ("The team went from two to one, and the protocol followed.",
     "The standardised pendulum ladder was displaced by an ambient-condition experiment \\u2014 "
     "a running dehumidifier, switched on and off \\u2014 because it tests the project's actual "
     "claim rather than characterising the instrument."),
    ("We pre-registered.",
     "Metrics, windows, exclusion rules and trial counts are frozen in the repository, dated, "
     "before the data existed."),
], top=2.20, gap=1.20, fill=OLDGOLD)''')


# =============================================================================
# 6 · § 5 — the Reflection slide the deck does not have
# =============================================================================
ANCHOR_EITHER = '''# =============================================================================
# NEW — if the sensors are not good enough
# ============================================================================='''
sub(ANCHOR_EITHER, '''# =============================================================================
# NEW — what we got wrong
# =============================================================================
s = new_slide()
title(s, "What we got wrong", "three, and the third one is about our own product")
for i, (col, head, body) in enumerate([
    (DEEPGOLD, "The pilot could not have succeeded.",
     "Two trials per condition. An exact permutation test on n = 2 has a minimum "
     "attainable p of 0.167, so no arrangement of the data could have cleared 0.05. "
     "We designed a study whose result was fixed before collection \\u2014 and we did not "
     "notice until we tried to analyse it. Six per condition puts the floor at 0.001."),
    (DEEPGOLD, "There is no single best channel.",
     "The derived vibration channel beats the raw accelerometer on sensitivity, and "
     "loses to it on cross-device agreement \\u2014 each device runs its own 200 ms window "
     "clock, so two phones disagree about a windowed statistic more than about the "
     "motion underneath it. Sensitivity and comparability trade against each other."),
    (GOLD, "The app let us mislabel every session, silently.",
     "All six pilot sessions are stored as condition \\u201ccontrolled\\u201d, including both "
     "slams. We are building an instrument to record what nobody wrote down, and it "
     "accepted a whole dataset written down wrong without a word. Found by using it; "
     "left uncorrected on the record."),
]):
    yy = 2.20 + i * 1.62
    hexagon(s, 0.85, yy - 0.02, 0.58, str(i + 1), fill=col, size=14)
    text(s, 1.72, yy, 10.6, 0.38, head, size=19, bold=True, color=NAVY)
    text(s, 1.72, yy + 0.42, 10.6, 1.05, body, size=14, color=MUTED, line=1.22)
notes(s, "PLACEHOLDER")


''' + ANCHOR_EITHER)


# =============================================================================
# 7 · Summary retuned to close the Reflection section
# =============================================================================
sub('''text(s, 0.85, 2.35, 11.5, 3.4, [
    ("The recorder works, and the derived vibration channel measures a real physical event "
     "at 13 to 109 times its own noise floor.",
     {"size": 20, "color": WHITE, "space": 16}),
    ("The pilot is a feasibility result, not a reproducibility result. It has two trials per "
     "condition, one operator and one device.",
     {"size": 20, "color": WHITE, "space": 16}),
    ("One usability defect worth reporting: all six pilot sessions are stored with condition "
     "\\"controlled\\", including both slams. The app accepted the mislabelling without warning.",
     {"size": 20, "color": GOLD, "space": 16}),
    ("The pre-registered study runs this week.",
     {"size": 20, "color": WHITE}),
], line=1.3)''',
    '''text(s, 0.85, 2.35, 11.5, 3.4, [
    ("The recorder works. All six channels run, and the derived vibration channel measures "
     "a real physical event at 13 to 109 times its own noise floor.",
     {"size": 20, "color": WHITE, "space": 16}),
    ("This is a feasibility result, not a reproducibility result \\u2014 and that distinction is "
     "the honest version of what one term buys.",
     {"size": 20, "color": WHITE, "space": 16}),
    ("Whether a phone is a good enough instrument is still open. Whether it is a good enough "
     "recorder is not: that part works either way.",
     {"size": 20, "color": GOLD, "space": 16}),
    ("The pre-registered study runs this week.",
     {"size": 20, "color": WHITE}),
], line=1.3)''')


# =============================================================================
# 8 · dividers, then the declarative reorder
# =============================================================================
sub('''prs.save(OUT)''',
    '''# =============================================================================
# section dividers — authored last, positioned by ORDER
# =============================================================================
divider("1", "Aims and objectives",
        "what we set out to build, and the problem it is aimed at")
divider("2", "Project presentation",
        "what we built, and why consumer hardware makes it a question")
divider("3", "Changes to the plan",
        "four, and what each one cost or bought")
divider("4", "Results",
        "one pilot, three measured findings")
divider("5", "Reflection",
        "what we now believe, what we got wrong, and why either answer is useful")


# =============================================================================
# narration order — the single place slide sequence is decided
# =============================================================================
ORDER = [
    "__title__",
    "div:Aims and objectives",
    "idea0", "idea1", "idea2",
    "aims",
    "div:Project presentation",
    "Implementation",
    "Why this is hard",
    "In use",
    "Channels",
    "One test per sensor",
    "Cross-device alignment",
    "Privacy properties",
    "div:Changes to the plan",
    "Changes since the proposal",
    "div:Results",
    "Pilot study",
    "Derived vibration channel",
    "Effect of metric choice",
    "Sync fiducial recovery",
    "Unlabelled event detection",
    "div:Reflection",
    "Midway hypothesis",
    "What we got wrong",
    "Either outcome is useful",
    "Scope and limitations",
    "Standardising the disturbance",
    "Remaining work",
    "Contributions",
    "Summary",
]

REG["__title__"] = prs.slides[0]
missing = [k for k in ORDER if k not in REG]
extra = [k for k in REG if k not in ORDER]
if missing or extra:
    sys.exit(f"ORDER mismatch\\n  missing from REG: {missing}\\n  never ordered: {extra}")

_ids = {id(sl): e for sl, e in zip(prs.slides, list(sld_lst))}
for e in list(sld_lst):
    sld_lst.remove(e)
for k in ORDER:
    sld_lst.append(_ids[id(REG[k])])

prs.save(OUT)''')

p.write_text(s)
print("build_deck.py patched: dividers, reveals, two-tier implementation, reflection slide, ORDER pass")
