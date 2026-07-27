"""Speaker notes and the shown/hidden cut, keyed to the video grading rubric.

The rubric caps length at 8 minutes and calls the maximum strict, so every shown
slide carries a second budget and the total is asserted at the bottom — the build
fails rather than quietly running long. Notes are bullets to talk from, not a
script to read; reading a script on camera is audible.

Hiding rather than deleting (`show="0"`) keeps the deck the graders read and the
deck filmed as one artefact: the cut material is still there for the report.

Two facts are open at the time of writing and are marked [CONFIRM] inline —
whether light and micLevel have exported a session or are only streaming live,
and whether more than one device runs the dev client. Both are one-line swaps.
"""
import re
import sys

from pptx import Presentation

DECK = "Covariate_Demo.pptx"

HIDE = {10, 11, 12, 13, 14, 21, 22, 27, 28, 29}

NOTES = {

# -- section 0 ---------------------------------------------------------------
1: """≈15s  ·  OPEN

• Covariate. A feasibility study: can consumer-grade phone sensors capture ambient
  experimental context well enough to be worth attaching to a record?
• CS-7470, Team 42. Caitlin Everett and Christopher Kimberley.
• Flat statement of what the next seven minutes contain — aims, what we built,
  what changed, what we measured, what we think now.
""",

# -- section 1 · AIMS AND OBJECTIVES -----------------------------------------
2: """≈3s  ·  DIVIDER

• Say the section name out loud: "Aims and objectives."
""",

3: """≈14s  ·  WHY RECORD THE ROOM  (1 of 3 — no papers yet)

• Two people run the same written protocol and get different answers.
• The difference usually isn't in the protocol. It's in everything around it that
  nobody wrote down.
• Bread is the everyday version — same recipe, different kitchen, different loaf.
  Nobody writes down the humidity.
""",

4: """≈12s  ·  + COLLBERG & PROEBSTING  (2 of 3)

• Start in our own field, because it's not somebody else's problem. Collberg and
  Proebsting went after 601 papers from eight ACM conferences and five journals.
• They tried to get the code and build it. Succeeded for 32.3%.
• Most published computer systems work cannot be repeated from what was written
  down — and that's the field with version control.
""",

5: """≈12s  ·  + THE ISOS CONSENSUS  (3 of 3)

• The interesting case is what a field does once it admits this.
• Perovskite solar cells: results weren't comparable between labs, because
  everybody tested under different ambient conditions and reported different
  parameters. Nature Energy, 2020 — a consensus statement, dozens of groups.
• The fix wasn't a better measurement. It was agreeing on which conditions have to
  be written down: temperature, humidity, illumination, atmosphere, electrical bias.
• That's the whole argument in one citation. The ambient conditions were always
  affecting the result. They just weren't in the record until somebody decided they
  had to be.
• So: every room where an experiment happens already contains a capable sensor
  package in somebody's pocket. The question is whether what it measures is worth
  anything.
""",

6: """≈38s  ·  AIMS AND OBJECTIVES  (as proposed)

• Three objectives, stated as we proposed them.
• BUILD — an ambient-context recorder. Pressure, motion, magnetic field, light,
  sound level, all on one clock, exported as a file bound to a named experiment.
• REPRODUCE — two known sensing techniques on commodity hardware, and see what
  survives the move off instrument-grade sensors.
• EVALUATE — whether any of it is trustworthy. Sampling health, agreement between
  devices, and whether a logged covariate actually explains variation between runs.
• The constraint we set ourselves: no additional hardware. It has to run on a phone
  a lab already owns, or nobody uses it.
• Worth saying plainly — this was always framed as a feasibility study. Our TA held
  us to that and was right to.
""",

# -- section 2 · PROJECT PRESENTATION ----------------------------------------
7: """≈3s  ·  DIVIDER

• "Project presentation" — what we built.
""",

8: """≈45s  ·  IMPLEMENTATION

• React Native. Four kinds of channel.
• DIRECT — straight out of expo-sensors. Accelerometer 50 Hz, magnetometer 25,
  barometer event-driven at about 1.
• DERIVED — vibration. The one piece of real engineering here. Gravity is a
  constant 1 g; a door closing moves the accelerometer about 1%. Estimate gravity
  with a low-pass, subtract it, then RMS and peak over a 200 ms window. Same
  hardware, different question, several times the signal-to-noise.
• NATIVE — light and microphone level. Swift on iOS, Kotlin on Android, because
  Expo exposes neither. Sound is stored as a LEVEL. No audio is recorded, so there
  is no waveform to leak — enforced by the code, not by a policy.
• RECORD — one session, one JSON file. Metadata, placement, per-channel sampling
  health, one shared clock. Two devices align on a haptic fiducial, because the
  timestamps are per-device and share no origin.
• THE POINT OF THE SLIDE — there's a deployment boundary through the middle. Four
  channels install by scanning a QR code: no build, no developer account, no cable.
  All six need a compiled dev client.
• [CONFIRM] All six now run on my device. Say "running" if light and mic haven't
  exported a session yet; say "recording" once they have.
• That gap is a finding, not an excuse. The hardware is universal. Access to it
  isn't. For a project whose premise is "everyone already carries this," that's
  the interesting part.
""",

9: """≈35s  ·  WHY THIS IS HARD

• Why this is a study and not a product announcement. Consumer sensors have
  documented, named failure modes.
• The raw stream isn't raw — SensorID, IEEE S&P 2019, recovers factory calibration
  baked into firmware. You're reading a vendor-conditioned number, not a physical
  quantity.
• Self-heating — a MEMS gyro drifts 317 degrees an hour in its first 400 seconds
  from power-on alone. A phone logging continuously starts by measuring itself.
  We think that's what our barometer warm-up is.
• Per-unit thermal drift — four identical units measured from −1.2 to +1.4 mg per
  degree against a ±0.5 spec. No generic correction exists; every device needs its
  own.
• Automatic pipelines — auto-exposure and automatic gain control sit between the
  world and the number. Both have to be defeated before light or sound means
  anything.
• Device heterogeneity — you can predict which OS a device runs from its sensor
  quality metrics alone, at 98% accuracy. A multi-device study measures devices as
  much as it measures rooms.
• No traceability — no calibration chain back to a national standard. So: relative
  change within a session. Never an absolute value.
• Plus two of our own — occlusion, meaning a fingerprint on the lens, lint in the
  mic port, a case over the barometer vent; and the OS suspending background apps
  on long overnight runs.
""",

10: """[hidden]  IN USE — recording screen, live values, sampling health, the O1 gate.""",
11: """[hidden]  CHANNELS — per-channel rates and analysis approach.""",
12: """[hidden]  ONE TEST PER SENSOR — the decisive experiment for each channel.""",
13: """[hidden]  CROSS-DEVICE ALIGNMENT — folded into Implementation.""",
14: """[hidden]  PRIVACY PROPERTIES — folded into Implementation.""",

# -- section 3 · CHANGES TO THE PLAN -----------------------------------------
15: """≈3s  ·  DIVIDER

• "Changes to the plan."
""",

16: """≈42s  ·  CHANGES SINCE THE PROPOSAL

• Four of them.
• ONE — scope was cut to four channels, then won back to six. Light and sound need
  a compiled dev client, so the pilot ran without them and we dropped the
  Alka-Seltzer study. The dev client now builds and all six run. What looked like
  a dropped feature turned into a measured deployment boundary, which is a better
  outcome than the original plan would have given us.
• TWO — the multi-site study is now explicitly a case study. Our TA pointed out
  that three sites with one person each can't support a claim about between-site
  variance: person, city, phone model and building are completely confounded.
  That's correct. The quantitative claims moved to a within-site design.
• THREE — the team went from two to one, and the protocol followed. The
  standardised pendulum ladder got displaced by an ambient-condition experiment —
  a loud dehumidifier, switched on and off, with an identical door-close protocol
  under both. That tests the project's actual claim, which is whether an
  unrecorded ambient condition changes the measured outcome, rather than just
  characterising the instrument.
• FOUR — we pre-registered. Metrics, windows, exclusion rules and trial counts,
  frozen in the repo, dated, before the data existed.
""",

# -- section 4 · RESULTS -----------------------------------------------------
17: """≈3s  ·  DIVIDER

• "Results."
""",

18: """≈45s  ·  PILOT STUDY   [footage runs here]

• The pilot: two baselines, two normal door closes, two slams. One phone, one room,
  Toronto. Chris ran it and narrates it.
• [CUE] Let the footage run. Caption over any cut: "protocol continues — 4 further
  trials."
• Coming out of the footage, the headline: the derived vibration channel separated
  slams from baseline at 13 to 109 times the session noise floor.
• And a correction to our own pilot report. It said the sync taps were lost. They
  weren't — re-running the detector at 50 Hz found 3 to 5 taps in all six sessions.
  A reporting artifact, not a data loss. Worth saying because we wrote the wrong
  thing down first.
• THIS SLIDE IS THE TIME LEVER. If the run is going long, shorten the footage here.
""",

19: """≈30s  ·  DERIVED VIBRATION CHANNEL

• The engineering result. Raw accelerometer against the derived channel, same
  events.
• Gravity dominates: a door close is about 1% on top of a constant 1 g. Subtract
  the constant and the same samples give 3.5 to 5 times the signal-to-noise. No new
  hardware.
• One trap worth naming, because we walked into it. You cannot take a spectrum of
  the acceleration MAGNITUDE. Magnitude is the square root of signal-squared plus
  one, which approximates to one plus signal-squared over two — so a transverse
  vibration enters squared and comes out at double its real frequency. We planted a
  7 Hz signal and the magnitude spectrum returned 14. Per-axis and detrended, it
  returns 7.03.
• That's the kind of mistake that produces a confident wrong number, so it's in the
  report.
""",

20: """≈25s  ·  EFFECT OF METRIC CHOICE

• Same four events, three metrics — window energy, peak, RMS.
• All three order the conditions correctly. What differs is the margin: 1.38 times,
  1.94, and 2.70.
• So metric choice doesn't change the direction here. It changes how much room you
  have before noise eats the result. That's why the metric is frozen in the
  pre-registration instead of picked after looking at the data.
• Honest note: we first wrote that one of the three failed to separate them. Then
  we rendered the figure and it hadn't. We fixed the claim.
""",

21: """[hidden]  SYNC FIDUCIAL RECOVERY — covered verbally on the pilot slide.""",
22: """[hidden]  UNLABELLED EVENT DETECTION — in the report.""",

# -- section 5 · REFLECTION --------------------------------------------------
23: """≈3s  ·  DIVIDER

• "Reflection." Say it — this is where the marks are.
""",

24: """≈36s  ·  MIDWAY HYPOTHESIS

• Rather than wait for the paper, here's what we believe right now per channel, and
  what would change our minds. Stated before the data, so it can be wrong.
• VIBRATION — likely good enough. 13 to 109 times the floor, about 1% repeatability
  within a condition. Falsified if the pendulum dose ladder isn't monotonic.
• ACCELEROMETER — good for timing, not for amplitude. It recovered every fiducial
  and every event, but events sit only 1 to 4% above gravity. Falsified if fiducials
  fail to align two devices.
• BAROMETER — not yet trusted. The session-long pressure rise looked identical in
  baseline and in slam, which says we're measuring the device, not the room. It
  earns trust if it tracks a weather station over eight hours.
• MAGNETOMETER — probably not useful at this scale. Event deviation 0.31 to 0.56
  microtesla against a baseline spread of 1.03. The noise is bigger than the signal.
  A clean negative is still a result.
• [CONFIRM] LIGHT and MIC LEVEL — now running. Position is untested but no longer
  blocked. Update this line once tonight's runs are in.
• One probably-good, one useful-for-timing, one suspect, one probably-not, two open.
  Naming which is which is most of what a feasibility study is for.
""",

25: """≈35s  ·  WHAT WE GOT WRONG

• Three, and the third one is about our own product.
• ONE — the pilot could not have succeeded. Two trials per condition. An exact
  permutation test on n equals two has a minimum attainable p of 0.167, so no
  arrangement of that data could ever have cleared 0.05. We designed a study whose
  result was fixed before we collected anything, and we didn't notice until we tried
  to analyse it. Six per condition puts the floor at 0.001 — which is why the
  pre-registration specifies six.
• TWO — there is no single best channel. The derived channel wins on sensitivity and
  loses on cross-device agreement, because each device runs its own 200-millisecond
  window clock. Two phones disagree about a windowed statistic more than they
  disagree about the motion underneath it. Sensitivity and comparability trade
  against each other, and we didn't expect that.
• THREE — the app let us mislabel every session, silently. All six pilot sessions
  are stored with condition "controlled". Including both slams. We're building an
  instrument to record what nobody wrote down, and it accepted a whole dataset
  written down wrong without a word of complaint. We found it by using our own tool,
  and we've left it uncorrected on the record.
""",

26: """≈20s  ·  EITHER OUTCOME IS USEFUL

• Close on why the question is worth asking regardless of which way it goes.
• IF THE SENSORS PASS — ambient context capture is free for anyone with a phone.
  Teaching labs, citizen science, anywhere with no instrumentation budget.
• IF THEY FAIL — the app is still the experiment-linked recorder: the schema, the
  UI, the sampling-health gate. The sensing moves to a cheap external package over
  BLE, and the phone's own sensors stop being confounding variables.
• That second case has support. The same phone sound-level apps, re-run with
  external calibrated microphones, came within ±1 dB of reference. The built-in
  signal chain was the limit, not the phone.
• So the real question isn't "does the phone work." It's "what's the cheapest thing
  that does, and can this app host it."
""",

27: """[hidden]  SCOPE AND LIMITATIONS — distributed into the two reflection slides.""",
28: """[hidden]  STANDARDISING THE DISTURBANCE — superseded by the ambient experiment.""",
29: """[hidden]  REMAINING WORK — in the report.""",

30: """≈12s  ·  CONTRIBUTIONS

• Chris ran the door pilot — protocol, recording, the six sessions.
• The recorder, export schema, analysis and study design are mine.
""",

31: """≈15s  ·  SUMMARY

• The recorder works. All six channels run.
• This is a feasibility result, not a reproducibility result — and that distinction
  is the honest version of what one term buys.
• Whether a phone is a good enough INSTRUMENT is still open. Whether it's a good
  enough RECORDER isn't. That part works either way.
• The pre-registered study runs this week.
• [CUE] Cut back to the opening frame — phone on the counter, hold two seconds, out.
""",
}

# --- apply -------------------------------------------------------------------
prs = Presentation(DECK)
n = len(prs.slides._sldIdLst)
if set(NOTES) != set(range(1, n + 1)):
    sys.exit(f"notes cover {min(NOTES)}..{max(NOTES)} but the deck has {n} slides")

for i, slide in enumerate(prs.slides, 1):
    slide.notes_slide.notes_text_frame.text = NOTES[i].strip()
    slide._element.set("show", "0" if i in HIDE else "1")

# --- budget ------------------------------------------------------------------
total = 0
for i, body in NOTES.items():
    if i in HIDE:
        continue
    m = re.search(r"≈(\d+)s", body)
    if not m:
        sys.exit(f"shown slide {i} carries no budget marker")
    total += int(m.group(1))

prs.save(DECK)
shown = n - len(HIDE)
print(f"{n} slides · {shown} shown · {len(HIDE)} hidden {sorted(HIDE)}")
print(f"budget {total}s = {total // 60}:{total % 60:02d}   cap 8:00, margin {480 - total}s")
if total > 470:
    sys.exit("OVER BUDGET — trim before filming")
