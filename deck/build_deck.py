"""Build the Covariate demo deck on the Georgia Tech 2022 template.

Design decisions worth stating, because they are the ones a reader would query:

* **Brand colours for furniture, validated palette inside the figures.** GT navy
  (#003057) and Tech gold (#EAAA00) carry titles, hexes and rules. They do NOT
  carry data series: run through the categorical-palette checks, navy and GT teal
  fall below the chroma floor (they read gray) and Tech gold falls outside the
  lightness band and under 3:1 on white. Brand colours are chosen for a logo on a
  building, not for telling four lines apart at 12pt. The figures keep the
  already-validated blue/orange/aqua.
* **The hexagon is the motif**, taken from the template's own shape language
  rather than invented — GT's deck already uses hex clusters. It numbers list
  items and holds stat callouts, which keeps the deck recognisably Georgia Tech
  and not recognisably generic.
* **Dark navy for opening, section and closing slides; white for content.** The
  sandwich, so the deck has a shape rather than eighteen identical pages.
* No accent stripes, no rules under titles.

Speaker notes carry the full narration, so the deck can be filmed straight
through. `[CUE]` lines mark where footage is cut in.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE, OUT = "gt.pptx", "Covariate_Demo_Full.pptx"
FIG = Path("figures")

NAVY = RGBColor(0x00, 0x30, 0x57)
GOLD = RGBColor(0xEA, 0xAA, 0x00)
OLDGOLD = RGBColor(0xB3, 0xA3, 0x69)
TEAL = RGBColor(0x00, 0x8C, 0x95)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x54, 0x58, 0x59)
# GT gold measures about 2:1 on white. It fills shapes; it never sets type on a
# light surface. DEEPGOLD is the darkened step for the cases where a gold-family
# text colour is genuinely wanted.
DEEPGOLD = RGBColor(0x7A, 0x5C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF2, 0xF3, 0xF4)

W, H = 13.333, 7.5
FONT = "Arial"

prs = Presentation(TEMPLATE)
# Drop the template's own 22 example slides; keep its layouts, theme and logo.
sld_lst = prs.slides._sldIdLst
for sld in list(sld_lst):
    # Drop the relationship as well as the list entry. Removing only the entry
    # leaves the slide parts in the package: python-pptx then reuses their
    # numbers for new slides, and the saved zip carries two members called
    # slide18.xml. PowerPoint reads that as a corrupt file.
    prs.part.drop_rel(sld.rId)
    sld_lst.remove(sld)

PLAIN = prs.slide_layouts[2]      # Title - Plain: GT logo, subtle background
TITLE_TT = prs.slide_layouts[3]   # Title - Tech Tower


# --- helpers -----------------------------------------------------------------


def _clear_placeholders(slide):
    """Template title layouts carry two placeholders; we position our own text."""
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def new_slide(layout=PLAIN, dark=False):
    s = prs.slides.add_slide(layout)
    _clear_placeholders(s)
    if dark:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
        bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background(); bg.shadow.inherit = False
        s.shapes._spTree.remove(bg._element)
        s.shapes._spTree.insert(2, bg._element)   # behind everything, above bg art
    return s


def text(slide, x, y, w, h, runs, size=16, color=INK, bold=False,
         align=PP_ALIGN.LEFT, space=6, line=1.0, italic=False):
    _italic = italic
    """runs: str, or list of (text, {overrides}) tuples, one paragraph each."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        body, over = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.space_after = Pt(over.get("space", space))
        p.line_spacing = over.get("line", line)
        r = p.add_run(); r.text = body
        f = r.font
        f.name = FONT
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        f.italic = over.get("italic", _italic)
    return tb


def hexagon(slide, x, y, side, label, fill=NAVY, fg=WHITE, size=15):
    """The motif. Small GT hex holding a numeral, letter or short stat."""
    sh = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y),
                                Inches(side), Inches(side * 0.88))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = fg
    return sh


REG: "dict[str, object]" = {}


def register(slide, key):
    if key in REG:
        sys.exit(f"duplicate slide key: {key}")
    REG[key] = slide
    return slide


def title(slide, t, sub=None, dark=False, k=None):
    text(slide, 0.85, 0.62, 11.6, 1.0, t, size=36, bold=True,
         color=WHITE if dark else NAVY)
    if sub:
        text(slide, 0.85, 1.42, 11.6, 0.5, sub, size=16,
             color=RGBColor(0xC8, 0xD2, 0xDC) if dark else MUTED)
    register(slide, k or t)


def divider(num, name, answers):
    """A dark section marker naming a rubric line in the grader's own words."""
    d = new_slide(dark=True)
    hexagon(d, 0.85, 2.60, 1.05, num, fill=GOLD, fg=NAVY, size=30)
    text(d, 2.35, 2.62, 10.2, 0.9, name, size=42, bold=True, color=WHITE)
    text(d, 2.35, 3.62, 10.2, 0.6, answers, size=17,
         color=RGBColor(0xC8, 0xD2, 0xDC), line=1.3)
    register(d, f"div:{name}")
    return d


def citation_card(slide, y, venue, head, authors, quote, h=1.74):
    """A paper reference as native shapes — citation plus the exact sentence.

    Not a page screenshot. A two-column journal page rendered into a 5-inch box is
    unreadable from a video, and the quoted sentence is the entire reason the paper
    is on the slide.
    """
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85),
                                   Inches(y), Inches(11.6), Inches(h))
    panel.fill.solid(); panel.fill.fore_color.rgb = PALE
    panel.line.color.rgb = OLDGOLD; panel.line.width = Pt(1.25)
    panel.shadow.inherit = False
    panel.text_frame.text = ""

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85),
                                  Inches(y), Inches(0.09), Inches(h))
    rule.fill.solid(); rule.fill.fore_color.rgb = GOLD
    rule.line.fill.background(); rule.shadow.inherit = False

    text(slide, 1.28, y + 0.18, 10.6, 0.26, venue, size=11.5, bold=True,
         color=DEEPGOLD)
    text(slide, 1.28, y + 0.46, 10.6, 0.32, head, size=16, bold=True, color=NAVY)
    text(slide, 1.28, y + 0.76, 10.6, 0.26, authors, size=12, color=MUTED)
    text(slide, 1.28, y + 1.06, 10.6, 0.60, quote, size=14, color=NAVY,
         italic=True, line=1.22)
    return panel


def hex_rows(slide, items, top=2.25, gap=1.18, x=0.85, wide=11.4, fill=NAVY):
    """Numbered hex + bold head + body. The deck's workhorse layout."""
    y = top
    for i, (head, body) in enumerate(items, 1):
        hexagon(slide, x, y - 0.04, 0.62, str(i), fill=fill)
        text(slide, x + 0.92, y, wide - 0.92, 0.4, head, size=19, bold=True, color=NAVY)
        text(slide, x + 0.92, y + 0.38, wide - 0.92, 0.7, body, size=14.5,
             color=MUTED, line=1.25)
        y += gap


def table(slide, cols, rows, top=2.30, row_h=0.62, size=13.5, head_size=12):
    """Simple column-positioned table. Columns are (label, x, width) in inches."""
    for label, x, _ in cols:
        text(slide, x, top, 3.0, 0.3, label, size=head_size, color=MUTED)
    y = top + 0.42
    for r in rows:
        for (_, x, w), cell in zip(cols, r):
            body, colour, bold = (cell if isinstance(cell, tuple) else (cell, MUTED, False))
            text(slide, x, y, w, row_h, body, size=size, color=colour,
                 bold=bold, line=1.2)
        y += row_h
    return y


def picture(slide, name, top=2.05, bottom=6.60, max_w=11.3, x=None):
    """Scale a figure to fit its band and centre it.

    Setting width alone is what pushed three charts off the bottom of the slide
    and over the Georgia Tech logo: these figures run 1.7-2.1 in aspect, so a
    10.5 in width is 5-6 in tall against about 4.5 in of usable height. Fit to
    whichever dimension binds, then centre what is left over.
    """
    from PIL import Image
    path = FIG / name
    iw, ih = Image.open(path).size
    avail_h = bottom - top
    w = min(max_w, avail_h * iw / ih)
    h = w * ih / iw
    left = (W - w) / 2 if x is None else x
    return slide.shapes.add_picture(str(path), Inches(left),
                                    Inches(top + (avail_h - h) / 2), Inches(w))


def cue(slide, label, detail):
    """A footage marker: what to play here, sized to be obvious while filming."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.4),
                                 Inches(2.9), Inches(8.5), Inches(2.1))
    box.fill.solid(); box.fill.fore_color.rgb = PALE
    box.line.color.rgb = OLDGOLD; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    box.text_frame.text = ""
    text(slide, 2.9, 3.35, 7.5, 0.6, label, size=24, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
    text(slide, 2.9, 4.05, 7.5, 0.8, detail, size=14, color=MUTED,
         align=PP_ALIGN.CENTER, line=1.25)


def cite(slide, body, y=6.98):
    """A source line along the bottom. Held to 10.2 in so it clears the GT logo."""
    return text(slide, 0.85, y, 10.2, 0.30, body, size=10.5, color=MUTED, line=1.15)


def autoplay(slide, movie):
    """Switch an embedded movie from click-to-play to play-on-slide-load.

    add_movie() writes a p:video whose media node starts on
    `<p:cond delay="indefinite"/>` — PowerPoint's "wait for a trigger". Setting
    the delay to zero starts it with the slide. The condition is rewritten in
    place rather than the timing tree rebuilt, so the structure stays whatever
    python-pptx emits.
    """
    from pptx.oxml.ns import qn
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        return movie
    n = 0
    for cond in timing.iter(qn('p:cond')):
        if cond.get('delay') == 'indefinite':
            cond.set('delay', '0')
            cond.attrib.pop('evt', None)
            n += 1
    if not n:
        raise RuntimeError("autoplay: no indefinite start condition found")
    return movie


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s.strip()


# =============================================================================
# 1 — title
# =============================================================================
s = new_slide(TITLE_TT, dark=True)
_hero = FIG / "hero_title.jpg"
if _hero.exists():
    from PIL import Image as _Im
    _iw, _ih = _Im.open(_hero).size
    _h = 5.55
    _w = _h * _iw / _ih
    s.shapes.add_picture(str(_hero), Inches(W - 0.85 - _w), Inches(0.95), Inches(_w))
text(s, 0.85, 2.45, 8.2, 1.2, "Covariate", size=60, bold=True, color=WHITE)
text(s, 0.85, 3.60, 8.2, 1.0,
     "Can the sensors in an ordinary phone record the conditions an experiment ran "
     "under \u2014 well enough to be worth keeping?", size=21, color=GOLD, line=1.25)
_qr = FIG / "qr_repo_dark.png"
if _qr.exists():
    s.shapes.add_picture(str(_qr), Inches(0.85), Inches(5.52), Inches(1.15))
    text(s, 2.20, 5.69, 5.6, 0.6, [
        ("Open code and open data.", {"size": 14, "bold": True, "color": GOLD, "space": 3}),
        ("github.com/CaitlinEverett/ambient-recorder",
         {"size": 13, "color": RGBColor(0xC8, 0xD2, 0xDC)}),
    ], line=1.2)
text(s, 0.85, 4.72, 8.2, 1.0, [
    ("CS-7470  Mobile & Ubiquitous Computing  ·  Team 42", {"size": 15}),
    ("Caitlin Everett  ·  Christopher Kimberley", {"size": 15}),
], color=RGBColor(0xC8, 0xD2, 0xDC))
notes(s, """
[CUE — open on the phone lying flat on a quiet kitchen counter, hold 3 seconds]

This is a quiet room. Nothing is happening in it.

Except the refrigerator cycled twice, someone closed a door down the hall, and the
floor is still ringing from a truck outside. None of that gets written down. If an
experiment ran on this bench today and failed to reproduce tomorrow, none of it
would be in the notebook either.

Covariate is a phone app that records the room, so that when an experiment doesn't
reproduce, there is something to look at.
""")

# =============================================================================
# NEW — bottom line up front (3-step build)
# =============================================================================
# Successive slides rather than PowerPoint animations: python-pptx cannot write a
# timeline, and a slide build survives export to PDF and to video untouched.
# One line of body per point: four points under a picture strip leave ~0.62 in per
# row, and a wrapped second line runs straight into the next heading.
# One line of body per point. The title slide already asks the question, so this
# starts at what we did rather than restating it.
BLUF = [
    (NAVY, "What we did.",
     "Closed a door 24 times, gently then hard, while two phones recorded six channels."),
    (TEAL, "It worked.",
     "Every hard close peaked higher than every gentle one — 12.3 dB louder, 2.6× "
     "the vibration. No overlap."),
    (DEEPGOLD, "And there is a lot left.",
     "The phones disagree on absolute values, one channel returned nothing, and we have "
     "only tested a few devices."),
]

BW, BGAP = 3.20, 0.34
BX = (W - (3 * BW + 2 * BGAP)) / 2

for _n in (1, 2, 3):
    s = new_slide()
    title(s, "The short version", "three sentences, before any of the detail",
          k=f"bluf{_n}")
    for idx, name in enumerate(("bluf_app.jpg", "bluf_chris.jpg", "bluf_caitlin.jpg")):
        x = BX + idx * (BW + BGAP)
        s.shapes.add_picture(str(FIG / name), Inches(x), Inches(1.88), Inches(BW))
    for idx, cap in enumerate(("six channels, live, on one clock",
                               "the pilot rig, Toronto",
                               "the same door protocol, Chicago")):
        text(s, BX + idx * (BW + BGAP), 4.06, BW, 0.26, cap, size=11.5,
             color=MUTED, align=PP_ALIGN.CENTER)
    # One line of body per point: 9.6 in keeps it clear of the Georgia Tech logo,
    # and a second wrapped line would run the third point into it.
    y = 4.52
    for k in range(_n):
        col, head, body = BLUF[k]
        hexagon(s, 0.85, y - 0.02, 0.52, str(k + 1), fill=col, size=13)
        text(s, 1.62, y - 0.03, 9.6, 0.32, head, size=18, bold=True, color=NAVY)
        text(s, 1.62, y + 0.31, 9.6, 0.30, body, size=13, color=MUTED)
        y += 0.64
    notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — the app, running, at slide size
# =============================================================================
# Full-bleed rather than a panel: the point of this slide is that a person can
# see the thing work, and a 2-inch phone on a white field does not carry that.
# The source is portrait, so the 16:9 frame is filled with a blurred, darkened
# copy of the same footage — the phone keeps its own aspect ratio and no part of
# the screen is cropped away.
s = new_slide(dark=True)
_APP = Path("media/chris_app.mp4")
if _APP.exists():
    autoplay(s, s.shapes.add_movie(str(_APP), Inches(0), Inches(0),
                                   Inches(W), Inches(H),
                                   poster_frame_image="media/chris_app_poster.jpg",
                                   mime_type="video/mp4"))
text(s, 0.85, 0.62, 6.4, 0.5, "Chris, running the app", size=30, bold=True,
     color=WHITE)
# Kept short so the line stays inside the dark half of the frame — the blurred
# fill is bright on the right and gold at 15 pt does not survive it.
text(s, 0.85, 1.20, 3.4, 0.9, "six channels reading live",
     size=15, color=GOLD, line=1.2)
register(s, "chris_app")

# =============================================================================
# 2 — the idea
# =============================================================================
# The same slide three times, adding one paper each. python-pptx cannot write
# PowerPoint animations, so the reveal is built as successive slides — which also
# survives export to PDF and to video without anyone clicking anything.
EDGE = dict(
    venue="Biochemistry and Biophysics Reports 26:100987 (2021)  \u00b7  "
          "doi.org/10.1016/j.bbrep.2021.100987",
    head="The edge effect: the trouble with culturing cells in 96-well plates",
    authors="Mansoury, Hamed, Karmustaji, Al Hannan & Safrany",
    quote="Wells around the rim of a plate read up to 35% lower than wells in the middle "
          "\u2014 same cells, same protocol, same plate. The cause is evaporation and a "
          "temperature gradient across the room.",
)

ISOS = dict(
    venue="Nature Energy 5:35\u201349 (2020)  \u00b7  doi.org/10.1038/s41560-019-0529-5",
    head="Consensus statement for stability assessment and reporting for perovskite "
         "photovoltaics",
    authors="Khenkin, Katz, Abate et al.  \u00b7  a whole field agreeing what to write down",
    quote="\u201cPublications lack consistency in experimental procedures and parameters "
          "reported. It is therefore challenging to reproduce and compare results.\u201d",
)

for _step in (0, 1, 2):
    s = new_slide()
    title(s, "Why record the room",
          "the problem this is aimed at, and why it belongs in a ubicomp course",
          k=f"idea{_step}")
    text(s, 0.85, 1.88, 11.4, 0.95, [
        ("Two experiments run to the same protocol still disagree, and the difference "
         "is the room.",
         {"size": 18, "bold": True, "color": NAVY, "space": 7}),
        ("Some labs do record it: ISO/IEC 17025 \u00a76.3.3 requires accredited labs to "
         "monitor and record environmental conditions, with validated per-room systems "
         "and scheduled recalibration. Everyone without a compliance budget writes down "
         "nothing.", {"size": 13, "color": MUTED}),
    ], line=1.25)
    # Two lines of intro at 18 pt end near 2.55 in. Cards stack from there with a
    # 0.12 in gutter; anything past 6.7 in collides with the Georgia Tech logo.
    if _step >= 1:
        citation_card(s, 3.02, **EDGE)
    if _step >= 2:
        citation_card(s, 4.88, **ISOS)
    notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — why consumer hardware makes this hard
# =============================================================================
s = new_slide()
title(s, "Why this is hard", "documented failure modes of consumer-grade sensors")
table(s,
      [("failure mode", 0.85, 3.5), ("evidence", 4.5, 4.4), ("consequence for us", 9.1, 3.4)],
      [(("The raw stream isn't raw", NAVY, True),
        "SensorID recovers factory calibration\nbaked into firmware (IEEE S&P 2019)",
        "we read a vendor-conditioned value,\nnot a physical quantity"),
       (("Self-heating", NAVY, True),
        "MEMS gyro drifts 317 °/h in its first\n400 s from power-on alone (2019)",
        "a phone logging continuously starts\nby measuring itself"),
       (("Per-unit thermal drift", NAVY, True),
        "four identical units: −1.2 to +1.4 mg/°C\nagainst a ±0.5 spec (2022)",
        "no generic correction exists;\nevery device needs its own"),
       (("Automatic pipelines", NAVY, True),
        "ambient-light cosine response off\n−33.87%; colorimetry needs locked ISO",
        "auto-exposure and AGC must be\ndefeated before light or sound counts"),
       (("Device heterogeneity", NAVY, True),
        "quality metrics alone predict the OS at\n0.98 accuracy (Sensors 2024)",
        "a multi-device study measures\ndevices as much as rooms"),
       (("No traceability", NAVY, True),
        "low-cost sensing lacks an unbroken\ncalibration chain to an NMI (2020)",
        "relative change within a session;\nnever an absolute value")],
      top=2.12, row_h=0.68, size=12.5)
text(s, 0.85, 6.72, 10.4, 0.4,
     "Two of our own to add: sensor occlusion, and the OS suspending background apps.",
     size=13, color=NAVY)
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — one decisive test per sensor
# =============================================================================
s = new_slide()
title(s, "One test per sensor", "each cheap, onsite, and sufficient to settle that channel")
table(s,
      [("channel", 0.85, 2.4), ("question", 3.3, 3.3), ("test", 6.7, 3.2),
       ("passes if", 10.0, 2.7)],
      [(("vibration / accel", NAVY, True), "tracks a graded mechanical dose?",
        "pendulum ladder, 5 levels × 6", "log-log slope CI excludes 0"),
       (("barometer", NAVY, True), "measuring air, or itself?",
        "8 h vs. a weather station,\nplus the warm-up curve", "slope ≈ 1; offset\nimplies true altitude"),
       (("magnetometer", NAVY, True), "usable at bench distance?",
        "magnet at 5 / 10 / 20 / 30 cm", "detectable at ≥ 10 cm;\nexponent near −3"),
       (("light", DEEPGOLD, True), "measurement or auto-exposure?",
        "lamp step, exposure locked\nand unlocked", "monotonic with lux,\nAE-independent"),
       (("micLevel", DEEPGOLD, True), "level, or AGC output?",
        "fixed tone, fixed distance,\nvarying background", "tracks the source,\nnot the background"),
       (("cross-device", NAVY, True), "do two models agree?",
        "same events, two devices,\none surface", "r ≥ 0.9; bias within\nthe noise floor")],
      top=2.12, row_h=0.68, size=12.5)
text(s, 0.85, 6.72, 10.4, 0.4,
     "Gold rows need a dev client. The rest run in one night plus about 75 minutes.",
     size=13, color=NAVY)
notes(s, "PLACEHOLDER")


# =============================================================================
# 3 — the plan
# =============================================================================
s = new_slide()
title(s, "Aims and objectives", "as proposed \u2014 three objectives and one design "
      "constraint", k="aims")
rows = [
    ("Build", "an ambient-context recorder: pressure, motion, magnetic field, light and sound "
              "level on one clock, exported as a file bound to a named experiment"),
    ("Reproduce", "two known sensing techniques on commodity hardware, and see what survives "
                  "the move off instrument-grade sensors"),
    ("Evaluate", "whether any of it is trustworthy — sampling health, agreement between "
                 "devices, and whether a logged covariate explains variation between runs"),
]
y = 2.35
for i, (head, body) in enumerate(rows):
    hexagon(s, 0.85, y - 0.05, 0.66, head[0], fill=[NAVY, TEAL, OLDGOLD][i])
    text(s, 1.85, y, 10.4, 0.4, head, size=21, bold=True, color=NAVY)
    text(s, 1.85, y + 0.42, 10.4, 0.7, body, size=15, color=MUTED, line=1.25)
    y += 1.30
text(s, 0.85, 6.30, 11.4, 0.8, [
    ("Design constraint: no additional hardware.", {"size": 17, "bold": True, "color": NAVY}),
    ("The recorder has to run on a phone a lab already owns.", {"size": 15, "color": MUTED}),
], space=3)
notes(s, """
The plan had three parts.

Build an ambient-context recorder — pressure, motion, magnetic field, light and
sound level, on one clock, exported as a file bound to a named experiment.

Reproduce two known sensing techniques on commodity hardware, and see what
survives the move off instrument-grade sensors.

And evaluate whether any of it is trustworthy: sampling health, agreement between
devices, and whether a logged covariate actually explains variation between runs.

The constraint was chosen deliberately — no special hardware. A phone every lab
already owns, or it doesn't get used.
""")

# =============================================================================
# 4 — what we built
# =============================================================================
s = new_slide()
title(s, "Implementation", "React Native; four channels install by QR code, six need "
      "a compiled build")
layers = [
    (NAVY, "1", "Direct sensors", "expo-sensors  \u2014  ships inside Expo Go",
     "accelerometer 50 Hz  \u00b7  magnetometer 25 Hz  \u00b7  barometer ~1 Hz"),
    (TEAL, "2", "Derived channel", "computed in JS from the raw stream",
     "vibration: gravity low-passed and subtracted, then RMS + peak over 200 ms"),
    (DEEPGOLD, "3", "Native modules", "Swift and Kotlin  \u2014  needs a dev build",
     "microphone LEVEL only, never audio  \u00b7  light via camera EXIF, because iOS "
     "exposes no ambient-light API to apps"),
    (OLDGOLD, "4", "Session record", "one session, one JSON file",
     "metadata  \u00b7  placement  \u00b7  per-channel sampling health  \u00b7  one shared clock"),
]
y = 2.24
for col, num, name, how, what in layers:
    hexagon(s, 0.85, y - 0.02, 0.56, num, fill=col, size=13)
    text(s, 1.66, y, 3.3, 0.35, name, size=18, bold=True, color=NAVY)
    text(s, 5.05, y + 0.04, 7.4, 0.35, how, size=13, color=DEEPGOLD, bold=True)
    text(s, 1.66, y + 0.42, 10.6, 0.5, what, size=14, color=MUTED, line=1.2)
    y += 1.06

text(s, 0.85, 6.22, 10.4, 0.62, [
    ("Expo Go is a pre-built container with a fixed set of native modules.",
     {"size": 16, "bold": True, "color": NAVY, "space": 4}),
    ("Anything outside that set cannot be loaded into it \u2014 so two channels cost us "
     "the thirty-second install.", {"size": 14, "color": MUTED}),
], line=1.2)
cite(s, "expo-sensors ships Accelerometer, Gyroscope, Magnetometer, Barometer, "
        "DeviceMotion, Pedometer and LightSensor (Android only).")
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — what it looks like to use (4-step build: designed, then built)
# =============================================================================
# python-pptx cannot write PowerPoint animations, so the reveal is successive
# slides. Phone frames are 0.466 aspect: at any width wide enough to READ, the
# height overruns a 16:9 slide. They are therefore sized as evidence-that-it-is-
# real, and the payload — the live channel table — gets its own slide at 5x.
SCREENS = [
    ("phone_mockup.png", "designed", "the spec we built from"),
    ("phone_home.png", "built", "sensor check: all six present"),
    ("phone_rec.png", "recording", "00:06 \u2014 six channels live"),
    ("phone_export.png", "exported", "one JSON file, 182 KB"),
]
PW, GAP = 2.02, 0.42
ROW_X = (W - (4 * PW + 3 * GAP)) / 2

for _n in (1, 2, 3, 4):
    s = new_slide()
    title(s, "In use", "what we specified, and what it actually became", k=f"inuse{_n}")
    for idx, (fname, cap, sub) in enumerate(SCREENS[:_n]):
        x = ROW_X + idx * (PW + GAP)
        s.shapes.add_picture(str(FIG / fname), Inches(x), Inches(2.05), Inches(PW))
        text(s, x, 6.42, PW, 0.28, cap, size=15, bold=True,
             color=OLDGOLD if idx == 0 else NAVY, align=PP_ALIGN.CENTER)
        text(s, x - 0.22, 6.70, PW + 0.44, 0.44, sub, size=11.5, color=MUTED,
             align=PP_ALIGN.CENTER, line=1.15)
    notes(s, "PLACEHOLDER")


# =============================================================================
# 5 — channels
# =============================================================================
s = new_slide()
title(s, "Channels", "four have recorded data so far; light and micLevel need a dev build")
cols = [("channel", 0.85), ("values", 4.15), ("units", 7.25), ("rate", 9.05)]
for name, x in cols:
    text(s, x, 2.15, 3.2, 0.3, name, size=12.5, color=MUTED)
rows = [
    ("accelerometer", "[x, y, z]", "g", "50 Hz", NAVY),
    ("magnetometer", "[x, y, z]", "µT", "25 Hz", NAVY),
    ("barometer", "[pressure, rel. altitude]", "hPa, m", "event-driven", NAVY),
    ("vibration", "[rms, peak]", "g", "5 Hz  ·  derived", TEAL),
    ("light", "[brightness]", "EV", "5 Hz  ·  dev build", DEEPGOLD),
    ("micLevel", "[rms]", "dBFS", "10 Hz  ·  dev build", DEEPGOLD),
    ("sync", "[pulse, of]", "—", "on demand  ·  fiducial", TEAL),
]
y = 2.62
for ch, vals, units, rate, col in rows:
    text(s, 0.85, y, 3.2, 0.32, ch, size=15.5, bold=True, color=col)
    text(s, 4.15, y, 3.0, 0.32, vals, size=14, color=MUTED)
    text(s, 7.25, y, 1.7, 0.32, units, size=14, color=MUTED)
    text(s, 9.05, y, 3.5, 0.32, rate, size=14, color=MUTED)
    y += 0.55
text(s, 0.85, 6.62, 11.0, 0.4,
     "The microphone channel stores a level in dBFS. No audio is recorded, so no waveform "
     "is displayed.", size=15, color=NAVY, bold=True)
notes(s, """
The seven channels, with their units and rates.

Two of them — light and microphone level — say "dev build". Those are the native
modules, and they only run in a compiled dev client, not in Expo Go. The app
reports that honestly rather than pretending.

Vibration is marked derived, and sync is the alignment fiducial. Both are ours
rather than the platform's.

And the line at the bottom is a design rule, not a caption: sound is a level,
never a waveform, because there is no audio to draw one from. Drawing a squiggle
would misrepresent the privacy guarantee.
""")

# =============================================================================
# 6 — the sync fiducial (video cue)
# =============================================================================
s = new_slide()
title(s, "Cross-device alignment", "implemented; no dual-device recording collected yet")
cue(s, "▶  Footage: a recording session in the app",
    "show Mark sync firing — three haptic pulses, one second apart, recorded with audio — "
    "then the export and the JSON meta / health blocks")
text(s, 1.15, 5.55, 11.1, 1.2,
     "Session time is monotonic from each device's own recording start, so two phones share no "
     "clock origin. A timestamp alone therefore cannot align them. Firing the vibration motor "
     "produces an event that any phone resting on the same surface registers through its own "
     "accelerometer, giving one device a known emission time and the others a signal to "
     "correlate against.", size=15, color=MUTED, line=1.25)
notes(s, """
[CUE — play the Mark sync clip, with sound. Three buzzes, one second apart.]

Mark sync fires three haptic pulses a second apart and writes each one into the
record as it fires.

It's physical on purpose. Session time is monotonic from each phone's own
recording start, so two devices share no clock at all — an offset of tens of
seconds between two files is normal and means nothing. A button that only wrote a
timestamp would align nothing.

Driving the vibration motor makes an event that every phone on the same surface
hears through its own accelerometer. One device gets ground truth about when it
fired; the others get a signal to cross-correlate against.

The one-second spacing is load-bearing, and I'll show you why in a minute.
""")

# =============================================================================
# 7 — privacy
# =============================================================================
s = new_slide()
title(s, "Privacy properties", "enforced by the implementation rather than by policy")
items = [
    "Audio is never recorded. The microphone channel stores a level in dBFS.",
    "Reference video is audio-free by construction — the flag is invariantly false.",
    "Location, if enabled at all, is a reverse-geocoded region and an altitude.\nNever coordinates. Absent by default.",
    "Recording is session-scoped and user-initiated. There is no background collection.",
]
y = 2.35
for t in items:
    hexagon(s, 0.85, y - 0.02, 0.42, "✓", fill=TEAL, size=13)
    text(s, 1.55, y, 10.8, 0.7, t, size=15.5, color=MUTED, line=1.25)
    y += 0.92
text(s, 0.85, 6.15, 11.4, 0.9, [
    ("Ambient sensor records can still identify a household.",
     {"size": 17, "bold": True, "color": NAVY}),
    ("Occupancy, daily routine, per-device sensor bias and barometric floor level are all "
     "inferable. The report documents these and the controls applied.",
     {"size": 14.5, "color": MUTED}),
], space=3)
notes(s, """
Four privacy properties that are structural rather than promised.

Audio is never recorded — the microphone channel stores a level. Reference video
is audio-free by construction; the flag is invariantly false. Location, if it's
enabled at all, is a reverse-geocoded region and an altitude, never coordinates,
and it's absent by default. And recording is session-scoped and user-initiated —
there is no background collection.

But the honest line is the one at the bottom. Our proposal originally said this
data holds no personal content, and we retract that. A longitudinal ambient record
is data about a household. Vibration and light reveal occupancy and routine;
per-device sensor bias is a fingerprint; absolute pressure implies a floor of a
building. The report says all of that, along with what we do about it.
""")

# =============================================================================
# 8 — what changed
# =============================================================================
s = new_slide()
title(s, "Changes since the proposal", "each one forced by something we hit")
hex_rows(s, [
    ("Alka-Seltzer dissolution \u2192 door closes.",
     "The dissolution study needed light and sound level. Both are native modules, and Expo "
     "Go \u2014 the thing that makes the app installable by scanning a code \u2014 cannot load "
     "them. We swapped to an event the four Expo Go channels can actually see."),
    ("Then the dev client built, and all six came back.",
     "What looked like a dropped feature turned into a measured result: four channels "
     "install in thirty seconds, six need a toolchain. The hardware is universal; access "
     "to it is not."),
    ("Multi-site study \u2192 case study.",
     "Our TA pointed out that three sites with one person each confound person, city, "
     "phone model and building. That was correct, and the quantitative claims moved to a "
     "within-site design with a trial count able to carry them."),
    ("We started measuring the instrument, not just the room.",
     "Once we could read our own numbers, three defects surfaced in a week \u2014 a clock, a "
     "health check, and a magnet. Each one changed what we build next."),
], top=2.20, gap=1.24, fill=OLDGOLD)
notes(s, """
Three things changed.

First, a scoping decision we'd make again. Two of our five channels — light and
sound level — are native modules, so they only run in a compiled dev client, not
in Expo Go. Getting a second site recording through that extra build step was
friction we chose not to spend the schedule on, so we dropped the Alka-Seltzer
dissolution study and ran a door experiment that needs only what installs by
scanning a QR code.

That's worth naming as a finding rather than an inconvenience. The gap between
"works on my device" and "runs at another site" is exactly the constraint mobile
and ubiquitous systems live inside, and we spent a week of the schedule learning
it firsthand.

Second, our reviewer pointed out that three sites with one participant each can't
support a claim about between-site variance — person, city, phone model and
building are completely confounded. So the multi-site study is now explicitly a
case study, and the quantitative claims moved to a within-site design with enough
trials to carry them.

Third, we pre-registered. Metrics, windows, exclusion rules and trial counts are
frozen in the repository, dated, before the data existed.
""")

# =============================================================================
# NEW — tonight's protocol and the arithmetic behind n = 6
# =============================================================================
# Two ASCII cards rather than slide text: the alignment is the artefact, and a
# PowerPoint text box re-kerns on another machine and shears the box drawing.
s = new_slide()
title(s, "The experiment, line by line",
      "one continuous recording, three devices, ten minutes")
# One clip in the deck, not two: the footage of the run itself lives on the
# next slide, where it sits beside the collaborator's pilot and earns the time.
picture(s, "mockup_protocol.png", top=1.95, bottom=6.80, max_w=8.05, x=0.85)
text(s, 9.30, 2.06, 3.25, 4.4, [
    ("2 \u00d7 3.",
     {"size": 16, "bold": True, "color": NAVY, "space": 7}),
    ("Machine {off, on} \u00d7 door {none, normal, slam}. The dehumidifier is the "
     "ambient condition nobody writes in a methods section.",
     {"size": 12.5, "color": MUTED, "space": 11}),
    ("One recording, not six.",
     {"size": 16, "bold": True, "color": NAVY, "space": 7}),
    ("Every cell on one clock and one thermal state, so the machine is the only "
     "thing that changes.",
     {"size": 12.5, "color": MUTED, "space": 11}),
    ("What it can answer:",
     {"size": 16, "bold": True, "color": DEEPGOLD, "space": 7}),
    ("whether a running appliance pushes a normal close below detectability while "
     "a slam still gets through.",
     {"size": 12.5, "color": MUTED}),
], line=1.28)
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — why n = 6
# =============================================================================
s = new_slide()
title(s, "Designing it was harder than running it",
      "the trial count fixes the smallest p obtainable, before any data exists",
      k="Why n = 6")
picture(s, "mockup_why_n6.png", top=1.95, bottom=6.90, max_w=7.4, x=0.85)
text(s, 8.55, 2.10, 3.95, 4.3, [
    ("This is the pilot's real defect.",
     {"size": 17, "bold": True, "color": NAVY, "space": 10}),
    ("Not that the pilot found nothing \u2014 that it could not have found anything. "
     "With two trials per condition the smallest p an exact permutation test can "
     "return is 0.167, and the corrected bar is 0.025.",
     {"size": 13.5, "color": MUTED, "space": 12}),
    ("We had designed a study whose result was fixed before collection, and did "
     "not notice until we tried to analyse it.",
     {"size": 13.5, "bold": True, "color": DEEPGOLD, "space": 12}),
    ("The rule it gave us.",
     {"size": 17, "bold": True, "color": NAVY, "space": 10}),
    ("When the schedule tightens, cut conditions, not replicates. Replicates make "
     "a test capable of answering; conditions make it interesting.",
     {"size": 13.5, "color": MUTED}),
], line=1.28)
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — what the instrument measured
# =============================================================================
s = new_slide()
title(s, "What the instrument measured",
      "24 door trials, two devices recording simultaneously, one table")
table(s,
      [("", 0.85, 2.5), ("normal close  (n = 13)", 3.5, 3.3),
       ("slam  (n = 11)", 6.9, 3.0), ("separation", 10.0, 2.6)],
      [(("acoustic peak", NAVY, True), "\u221226.1 dBFS\n[\u221233.5, \u221221.8]",
        ("\u221213.8 dBFS\n[\u221218.3, \u221210.8]", TEAL, True), ("+12.3 dB\nno overlap", TEAL, True)),
       (("derived vibration", NAVY, True), "0.0039 g\n[0.0023, 0.0062]",
        ("0.0102 g\n[0.0071, 0.0125]", TEAL, True), ("2.60\u00d7\nno overlap", TEAL, True)),
       (("replicated on device 2", NAVY, True), "4.3\u00d7 noise floor",
        ("12.0\u00d7 noise floor", TEAL, True), ("2.77\u00d7\nno overlap", TEAL, True))],
      top=2.30, row_h=0.88, size=13.5)
text(s, 0.85, 5.55, 11.5, 1.1, [
    ("Not one normal close exceeded any slam \u2014 on either channel, on either device.",
     {"size": 18, "bold": True, "color": NAVY, "space": 7}),
    ("Complete separation at 13 versus 11 gives a minimum attainable one-tailed "
     "p of 1/C(24,13) = 4.0 \u00d7 10\u207b\u2077. This is the confirmatory dose\u2013response the "
     "analysis plan was built around, and it replicated independently on a second "
     "device in the same run.", {"size": 14, "color": MUTED}),
], line=1.3)
cite(s, "Noise floor and windowing follow CS-7470 L5-06 (Noise). Gravity-subtracted "
        "derivation after Mizell, ISWC 2003, and CS-7470 L5-03 (IMU / Gravity).")
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — cross-device agreement, and cross-device disagreement
# =============================================================================
s = new_slide()
title(s, "Two phones, one table",
      "they agree about change and disagree about absolute value \u2014 both are results")
for i, (col, head, body) in enumerate([
    (TEAL, "Agreement on change:  r = 0.97",
     "Two devices, 24 shared events. Pearson r = 0.970 on acoustic peak and 0.971 on "
     "the derived vibration channel, against a threshold we fixed in advance of 0.90. "
     "Event timing matched to a median 33 ms \u2014 one sample \u2014 with no fiducial, on "
     "starts 76 ms apart. H2 passes."),
    (DEEPGOLD, "Disagreement on absolutes:  measured, not cited",
     "Same table, same second: barometers differ by 0.675 hPa (5.6 m of equivalent "
     "altitude); resting acoustic level by 7.8 dB; noise floors across four devices by "
     "2.4\u00d7. One device read the magnetic field at 664 \u00b5T while the other read 41 \u2014 "
     "a magnet was attached to it, and no other channel noticed."),
]):
    yy = 2.20 + i * 2.05
    hexagon(s, 0.85, yy - 0.02, 0.60, "\u2713" if i == 0 else "\u0394", fill=col, size=15)
    text(s, 1.75, yy, 10.6, 0.4, head, size=20, bold=True, color=NAVY)
    text(s, 1.75, yy + 0.46, 10.6, 1.3, body, size=14, color=MUTED, line=1.3)
text(s, 0.85, 6.32, 10.4, 0.34,
     "Relative change within a session is now a measured constraint, not a borrowed "
     "caution.", size=15, bold=True, color=NAVY)
cite(s, "CS-7470 L5-07 (Sensor Calibration). Device-dependent bias: McNicholas & Mass, "
        "Wea. Forecasting 2021; Stisen et al., SenSys 2015; Kuhlmann et al., Behav. Res. "
        "Methods 2021.")
notes(s, "PLACEHOLDER")


# =============================================================================
# 9 — the pilot (video cue)
# =============================================================================
s = new_slide()
title(s, "Pilot study", "the same door protocol in two cities, eight weeks apart")
# Embedded rather than dragged in: build_deck.py regenerates the file from
# scratch, so anything added by hand in PowerPoint is lost on the next build.
# 16 s — Chris's Toronto trial, then the same protocol in Chicago.
MOVIE = Path("media/pilot_two_sites.mp4")
if MOVIE.exists():
    mv = s.shapes.add_movie(str(MOVIE), Inches(1.05), Inches(2.02),
                            Inches(2.60), Inches(4.62),
                            poster_frame_image="media/pilot_poster.jpg",
                            mime_type="video/mp4")
    autoplay(s, mv)
text(s, 4.35, 2.15, 8.0, 4.4, [
    ("Same protocol, two cities.",
     {"size": 20, "bold": True, "color": NAVY, "space": 9}),
    ("Christopher Kimberley ran the pilot in Toronto \u2014 two baselines, two normal door "
     "closes, two slams, on an iPhone X. The blue marker on the door edge is his "
     "repeatability control: the same closed position every trial.",
     {"size": 15, "color": MUTED, "space": 12}),
    ("The second half is the same action in Chicago, on different hardware, in a "
     "different building, eight weeks of app development later.",
     {"size": 15, "color": MUTED, "space": 12}),
    ("Six sessions, exported as JSON. Every result that follows was obtained from those "
     "files alone \u2014 in a different city, without further input from the person who "
     "recorded them.", {"size": 15, "color": MUTED, "space": 12}),
    ("That is the whole point of a portable record.",
     {"size": 15, "bold": True, "color": NAVY}),
], line=1.3)
notes(s, """
[CUE — play Chris's door-experiment footage, 10 to 15 seconds]

The pilot was run by Christopher Kimberley in Toronto: two baselines, two normal
door closes, two hard slams. One phone, one room, six sessions on an iPhone X.

Everything on the next three slides comes out of those six files. That's worth
noting on its own — the export format did its job. Somebody else's recordings, in
another city, reanalysed from scratch without asking him a single question.
""")

# =============================================================================
# 10 — the hack (fig2)
# =============================================================================
s = new_slide()
title(s, "Derived vibration channel", "signal-to-noise compared with the raw accelerometer")
picture(s, "deck_fig2_derived_vs_raw.png")
notes(s, """
Here's the derived channel, on a real slam.

The top trace is raw accelerometer magnitude. It peaks at 1.011 g — one percent
above gravity — because gravity is a constant one g sitting on top of everything
and the event is a rounding error next to it.

The bottom trace is the same sensor and the same samples, with a low-pass estimate
of gravity subtracted out. Twenty-eight times the noise floor.

Across the four door events the derived channel beats the raw sensor it comes from
by three-and-a-half to five times in signal-to-noise. That's the hack: subtracting
a large constant is what lets a small transient be seen at all.
""")

# =============================================================================
# 11 — metric choice (fig1)
# =============================================================================
s = new_slide()
title(s, "Effect of metric choice", "the same four events, measured three ways")
picture(s, "deck_fig1_metric_choice.png")
notes(s, """
Then a result about our own analysis.

These are the same four events measured three ways. The statistic the original
write-up used is a two-hundred-millisecond window average — and a
fifty-millisecond door impact gets diluted by wherever that window boundary
happens to fall.

All three orderings are correct. But the window average leaves a 1.4-times margin
between close and slam, where the energy integral leaves 2.7. Margin is what
survives more trials.

And with two trials per condition, none of it is significant — it could not have
been. The smallest p-value an exact permutation test can return at n equals two is
0.167. The design could not have reached significance on any data at all.

That's now a rule in the protocol: six trials per condition, minimum.
""")

# =============================================================================
# 12 — fiducial (fig3)
# =============================================================================
s = new_slide()
title(s, "Sync fiducial recovery", "the pilot report's missing taps were a reporting artifact")
picture(s, "deck_fig3_fiducial.png")
notes(s, """
One more, and this is the one I'd want a reviewer to see.

The pilot write-up concluded that only one of three sync taps had been recorded —
flagged as a data-quality problem.

They were all there. Three raps inside a few hundred milliseconds fall into one or
two windows of the five-hertz derived channel, which is where we looked. The
fifty-hertz raw accelerometer had them the whole time.

A reporting artifact, not a data failure — and exactly the kind of thing that gets
written into a paper if nobody re-runs the analysis. It's also why the in-app sync
marker spaces its pulses a full second apart.
""")

# =============================================================================
# 13 — blind detection (terminal cue)
# =============================================================================
s = new_slide()
title(s, "Unlabelled event detection", "run on the six pilot sessions; the detector receives no labels")
cue(s, "▶  Terminal: detect_events() on the pilot sessions",
    "18pt minimum")
text(s, 1.15, 5.45, 11.1, 1.3, [
    ("All four door events were recovered at the times the operator recorded. One baseline "
     "returned no candidates. The other returned a single candidate at 2.1 times the noise "
     "floor, which is either an unnoticed event or a false positive at the chosen threshold; "
     "it is reported as unresolved.", {"size": 15, "color": MUTED}),
], line=1.25)
notes(s, """
[CUE — terminal, run the blind detection live]

Last result, and it's the one that changes what we're allowed to claim.

Everything before this has the same shape: we caused an event, then we found it.
That demonstrates sensitivity and nothing about field performance, because the
analyst always knew where to look.

So: the detector gets no labels. It sets its own threshold from each recording's
quiet background and returns candidate events.

It finds all four door events, at the times Chris wrote down. It finds nothing in
one baseline — correctly. And it returns one marginal candidate in the other, at
about twice the noise floor, which is either a real event nobody noticed or a
false positive at our threshold. We report it either way.
""")

# =============================================================================
# NEW — midway hypothesis
# =============================================================================
s = new_slide()
title(s, "Midway hypothesis", "what we currently believe per channel, and what would change it")
table(s,
      [("channel", 0.85, 2.3), ("current position", 3.2, 3.1), ("on what evidence", 6.4, 3.2),
       ("would change if", 9.7, 3.0)],
      [(("vibration", TEAL, True), ("likely good enough", TEAL, True),
        "13–109× noise floor; ~1%\nwithin-condition repeatability",
        "the dose ladder is not\nmonotonic"),
       (("accelerometer", TEAL, True), ("good for timing,\nnot amplitude", TEAL, True),
        "recovered every fiducial and\nevent; 1–4% above gravity",
        "fiducials fail to align\ntwo devices"),
       (("barometer", DEEPGOLD, True), ("not yet trusted", DEEPGOLD, True),
        "session rise identical in\nbaseline and slam",
        "it tracks a weather\nstation over 8 h"),
       (("magnetometer", DEEPGOLD, True), ("probably not useful\nat this scale", DEEPGOLD, True),
        "event deviation 0.31–0.56 µT;\nbaseline spread 1.03 µT",
        "the magnet ladder is\ndetectable at 10 cm"),
       (("light", DEEPGOLD, True), ("sampling, but no\nvalue yet", DEEPGOLD, True),
        "87 samples in 6 s; the EV\nfield still renders \u2014",
        "an EV appears under\nknown illumination"),
       (("micLevel", TEAL, True), ("recording", TEAL, True),
        "\u221262.1 dBFS, 82 samples\nat 14 Hz, exported",
        "it tracks the source,\nnot the background"),
       (("cross-device", DEEPGOLD, True), ("one result already", DEEPGOLD, True),
        "realised rates differ 1.68\u00d7\nbetween our two phones",
        "the paired run returns\nr \u2265 0.9 on shared events")],
      top=2.05, row_h=0.60, size=12.5)
text(s, 0.85, 6.80, 10.4, 0.4,
     "Both blocked channels now record. One is already returning a result we did not "
     "predict.", size=13.5, color=NAVY)
notes(s, "PLACEHOLDER")


# =============================================================================
# 14 — limitations
# =============================================================================
s = new_slide()
title(s, "Scope and limitations", "a feasibility study, and what it does not establish")
text(s, 0.85, 2.15, 5.4, 0.35, "ESTABLISHED", size=15, bold=True, color=TEAL)
for i, t in enumerate(["Detects a real event at 13–109× its\nnoise floor",
                       "Repeatable to ~1% within a condition",
                       "Instrument characterised — warm-up,\ndrift, sampling health",
                       "A derived channel beats its raw\nsensor, 3–5×"]):
    text(s, 0.85, 2.74 + i * 0.80, 0.3, 0.3, "✓", size=14, bold=True, color=TEAL)
    text(s, 1.30, 2.74 + i * 0.80, 4.6, 0.7, t, size=14, color=MUTED, line=1.25)

text(s, 7.05, 2.15, 5.3, 0.35, "NOT ESTABLISHED", size=15, bold=True, color=DEEPGOLD)
for i, t in enumerate(["One participant per site — person, city,\nphone and building are confounded",
                       "One operator, one room, one device family",
                       "Light and sound level untested — dev build",
                       "Cross-device agreement — untested; one device"]):
    text(s, 7.05, 2.74 + i * 0.80, 0.3, 0.3, "—", size=14, bold=True, color=DEEPGOLD)
    text(s, 7.50, 2.74 + i * 0.80, 4.6, 0.7, t, size=14, color=MUTED, line=1.25)

text(s, 0.85, 5.75, 11.4, 1.2,
     "This follows our reviewer's assessment of the proposal: three sites with one participant "
     "each cannot demonstrate that logging a covariate reduces between-site variance. We accept "
     "that, and have moved the quantitative claims to a within-site design.",
     size=15, color=NAVY, line=1.3)
notes(s, """
What this work is, and what it isn't.

On the left, what it can claim. It's a feasibility study, and it succeeds as one:
a commodity phone detects a real physical event at thirteen to a hundred and nine
times its own noise floor, repeatably to about one percent within a condition. We
characterised the instrument itself — warm-up, drift, sampling health. And the
derived channel beats the raw sensor it comes from.

On the right, what it can't. One participant per site, so person, city, phone
model and building are all confounded. One operator, one room, one device family.
Light and sound level untested, because they need a dev build. Cross-device
agreement has one pair of devices behind it, at one site.

Our reviewer said it first: three sites with one participant each cannot show that
a logged covariate reduces between-site variance. They were right — so we moved
the quantitative claims to a within-site design and left the multi-site work as
what it is, a case study.
""")

# =============================================================================
# 15 — what we learned / the pendulum
# =============================================================================
s = new_slide()
title(s, "Standardising the disturbance", "designed, not yet run; two trials labelled slam differed by 3.9×")
picture(s, "deck_fig4_pendulum_design.png")
notes(s, """
[CUE — cut in the pendulum footage here: wide shot, then the close on the release]

The biggest lesson came from the slams. Two trials both labelled "slam" differed
by nearly a factor of four — comparable to the gap between slamming and closing.
"Hard" isn't a measurement, it's a mood, and we'd built a study on top of one.

So the disturbance gets a number. A fixed mass on a fixed string, released from
marked angles. Impact energy is m g L times one minus cosine theta — exact, and
five angles span a twenty-two-fold range.

Six trials at each of five levels, in randomised order. Six, because at two trials
per condition an exact test cannot return a p-value below 0.167.

What you're watching is a demonstration take — I'm talking next to the phone,
which our own frozen protocol excludes. The dataset run happens this week, in an
empty room.
""")

# =============================================================================
# 16 — queued
# =============================================================================
s = new_slide()
title(s, "Remaining work", "none of the following has been collected yet")
items = [
    ("Pendulum dose ladder", "5 levels × 6 trials, randomised order, empty room",
     "the confirmatory spine"),
    ("Overnight ambient runs", "8+ hours unattended, three nights",
     "the refrigerator duty cycle — a disturbance nobody notices"),
    ("Blind detection trial", "4 hours of ordinary activity, log sealed before analysis",
     "does it find events nobody labelled?"),
    ("Barometer vs. weather station", "12 hours against a National Weather Service record",
     "the only channel we can check outside our own project"),
]
y = 2.30
for i, (head, what, why) in enumerate(items, 1):
    hexagon(s, 0.85, y - 0.02, 0.52, str(i), fill=TEAL, size=13)
    text(s, 1.66, y, 10.6, 0.32, head, size=18, bold=True, color=NAVY)
    text(s, 1.66, y + 0.36, 10.6, 0.32, what, size=14, color=MUTED)
    text(s, 1.66, y + 0.70, 10.6, 0.32, why, size=14, color=TEAL,
         bold=False)
    y += 1.16
notes(s, """
What's queued for the rest of the week.

The pendulum dose ladder — five levels, six trials each, randomised order, empty
room. That's the confirmatory spine; everything quantitative rests on it.

Overnight ambient runs, eight hours or more, three nights. That's aimed at the
refrigerator duty cycle — a compressor cycling every half hour is precisely the
invisible variable we're arguing about, and nobody in the room notices it.

A blind detection trial: four hours of ordinary activity with the log sealed
before the analysis runs.

And a twelve-hour barometer run checked against a National Weather Service
station. That's the only channel we can validate against anything outside our own
project — and the residual offset should imply our own altitude, which we can
check against the known elevation of the room.
""")

# =============================================================================
# NEW — what we got wrong
# =============================================================================
s = new_slide()
title(s, "Things we are still working on",
      "three, and two of them are about our own product",
      k="What we got wrong")
for i, (col, head, body) in enumerate([
    (GOLD, "UI could guide people to make better choices about recording.",
     "Condition defaults to \u201ccontrolled\u201d, and site, notes and placement are all "
     "optional \u2014 so a session can be saved with none of them and the app never asks. "
     "Every session we have recorded, on both sides of the project, came out that way. "
     "Distance from phone to door is the largest single determinant of amplitude and it "
     "is in none of the files, which is why the two-site comparison can\u2019t be made "
     "rigorously. Nobody operated it wrong; the tool never made it easy to operate "
     "right. It is the first thing we are fixing."),
    (DEEPGOLD, "Six sensors appeared to disagree with their own spec.",
     "The dev client reported every channel overshooting its nominal rate by the same "
     "1.68\u00d7 \u2014 including a barometer reading 2 Hz when the hardware runs at 1. Six "
     "independent sensors do not agree on an error; a shared denominator does. The "
     "elapsed timer counted setInterval firings rather than seconds, and the JS thread "
     "was dropping them under sensor load."),
    (DEEPGOLD, "There is no single best channel.",
     "The derived vibration channel beats the raw accelerometer on sensitivity and loses "
     "to it on cross-device agreement, because each device runs its own window clock. "
     "Sensitivity and comparability trade against each other, and we did not expect "
     "that."),
]):
    yy = 2.12 + i * 1.66
    hexagon(s, 0.85, yy - 0.02, 0.58, str(i + 1), fill=col, size=14)
    text(s, 1.72, yy, 10.6, 0.38, head, size=18, bold=True, color=NAVY)
    text(s, 1.72, yy + 0.40, 10.6, 1.15, body, size=13, color=MUTED, line=1.2)
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — what happens next
# =============================================================================
s = new_slide()
title(s, "What happens next", "five days of it, and the version that outlives the term")
text(s, 0.85, 1.92, 5.6, 0.3, "THIS WEEK", size=12, bold=True, color=DEEPGOLD)
text(s, 6.95, 1.92, 5.5, 0.3, "AFTER THAT", size=12, bold=True, color=TEAL)
WEEK = [
    ("The other machine state", "completes the 2 \u00d7 3; ten minutes"),
    ("Magnet off, six closes", "turns tonight's staging error into a controlled test"),
    ("Log the light module", "returns nil, or throws? \u201cunsupported\u201d vs \u201cbroken\u201d"),
    ("The overnight run", "barometer against a weather station; the sampling gate"),
    ("Distance ladder \u2014 0.5, 1, 2 m", "makes Toronto and Chicago comparable at last"),
]
y = 2.32
for head, body in WEEK:
    text(s, 0.85, y, 5.6, 0.26, "\u2022  " + head, size=14, bold=True, color=NAVY)
    text(s, 1.16, y + 0.26, 5.3, 0.26, body, size=12.5, color=MUTED)
    y += 0.72

AFTER = [
    ("Baseline every device on the market",
     "every calibration run contributes a noise floor keyed to device model, so the app "
     "can tell any phone whether it resolves the effect you declared"),
    ("Refuse the devices that cannot",
     "a capability tier, not a disclaimer \u2014 warn or decline below threshold"),
    ("A channel-to-experiment matrix",
     "which sensors suit which kinds of experiment, published rather than guessed"),
    ("The sled",
     "external sensors over BLE where the phone's own are the limit"),
]
y = 2.32
for head, body in AFTER:
    text(s, 6.95, y, 5.5, 0.26, "\u2022  " + head, size=14, bold=True, color=NAVY)
    text(s, 7.26, y + 0.26, 5.2, 0.5, body, size=12.5, color=MUTED, line=1.2)
    y += 0.98

text(s, 0.85, 6.30, 10.4, 0.34,
     "The first column is homework. The second is the reason the homework is worth doing.",
     size=15, bold=True, color=NAVY)
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — if the sensors are not good enough
# =============================================================================
s = new_slide()
title(s, "Either outcome is useful", "the phone is the recorder; whether it is also the sensor is the question")
for i, (col, head, body) in enumerate([
    (TEAL, "If the sensors pass",
     "Ambient context capture becomes free for anyone with a phone — usable in "
     "teaching labs, in citizen science, and as opportunistic covariate logging where "
     "no instrumentation exists."),
    (DEEPGOLD, "If they fail",
     "The app remains the experiment-linked recorder, the schema and the UI, and the "
     "sensing moves to a cheap external package over BLE. Our ESP32 + BME280 stretch "
     "goal is already this architecture."),
]):
    x = 0.85 + i * 6.05
    hexagon(s, x, 2.30, 0.56, "✓" if i == 0 else "→", fill=col, size=14)
    text(s, x, 3.12, 5.5, 0.4, head, size=20, bold=True, color=NAVY)
    text(s, x, 3.62, 5.5, 1.6, body, size=14.5, color=MUTED, line=1.35)

# Held to 10.4 in wide and three lines: the fourth line reaches the GT logo.
text(s, 0.85, 5.62, 10.4, 1.3, [
    ("The second case has measured support.",
     {"size": 16, "bold": True, "color": NAVY, "space": 8}),
    ("The same sound-measurement apps, re-run with external calibrated microphones, came "
     "within ±1 dB of reference (Kardous & Shaw, JASA 2016). The built-in signal chain was "
     "the limit, not the phone.", {"size": 14.5, "color": MUTED}),
], line=1.3)
notes(s, "PLACEHOLDER")


# =============================================================================
# NEW — works cited
# =============================================================================
s = new_slide(dark=True)
title(s, "Works cited", dark=True)
# One paragraph per reference, wrapped by the renderer rather than by hand.
# The previous version hard-wrapped each entry across three list items with
# leading spaces for the hanging indent; any renderer that re-flows the box
# (Keynote, Google Slides, PowerPoint on iPad) collapsed those, which merged
# Khenkin into Mizell and ate a space in "Biochemistry and Biophysics".
LIT = [
    "Mansoury, Hamed, Karmustaji, Al Hannan & Safrany (2021). The edge effect: the "
    "trouble with culturing cells in 96-well plates. Biochemistry and Biophysics "
    "Reports 26, 100987. doi:10.1016/j.bbrep.2021.100987",

    "Khenkin, Katz, Abate et al. (2020). Consensus statement for stability assessment "
    "and reporting for perovskite photovoltaics. Nature Energy 5, 35\u201349. "
    "doi:10.1038/s41560-019-0529-5",

    "Mizell, D. (2003). Using gravity to estimate accelerometer orientation. "
    "ISWC \u201903, 252\u2013253.",

    "Zhang, Beresford & Sheret (2019). SensorID: sensor calibration fingerprinting for "
    "smartphones. IEEE Symposium on Security and Privacy.",

    "Stisen, Blunck, Bhattacharya et al. (2015). Smart devices are different. "
    "SenSys \u201915, 127\u2013140. doi:10.1145/2809695.2809718",

    "Peguero, Labrador & Cook (2016). Assessing jitter in sensor time series from "
    "Android mobile devices. IEEE SMARTCOMP. doi:10.1109/SMARTCOMP.2016.7501679",

    "McNicholas & Mass (2021). Bias correction of smartphone pressure observations. "
    "Weather and Forecasting 36(5), 1867\u20131889. doi:10.1175/WAF-D-20-0222.1",

    "Kuhlmann, Garaizar & Reips (2021). Smartphone sensor accuracy varies from device "
    "to device. Behavior Research Methods 53(1), 22\u201333. "
    "doi:10.3758/s13428-020-01404-5",

    "ISO/IEC 17025:2017, clause 6.3.3 \u2014 monitoring, control and recording of "
    "environmental conditions in accredited testing and calibration laboratories.",

    "Harris, ed. (2019). NISTIR 6969, Selected Laboratory and Measurement Practices "
    "and Procedures \u2014 laboratory siting and environmental requirements.",

    "Kardous & Shaw (2016). Evaluation of smartphone sound measurement applications "
    "using external microphones. JASA 140(4), EL327\u2013EL333. doi:10.1121/1.4964639",
]
COURSE = [
    "CS-7470 Mobile & Ubiquitous Computing \u2014 T. Ploetz, T. Starner",
    "",
    "L4-05   Magnetometer",
    "L4-06   Barometer",
    "L4-08   Ambient Light Sensor",
    "L5-03   IMU / Gravity",
    "L5-04, L5-05   Sampling Rate",
    "L5-06   Noise",
    "L5-07   Sensor Calibration",
]
text(s, 0.85, 2.16, 7.5, 4.7,
     [(r, {"size": 10, "color": RGBColor(0xC8, 0xD2, 0xDC), "space": 4})
      for r in LIT], line=1.18)
text(s, 8.75, 2.16, 3.8, 4.6, "\n".join(COURSE), size=11.5,
     color=RGBColor(0xC8, 0xD2, 0xDC), line=1.6)
text(s, 8.75, 1.80, 3.8, 0.3, "COURSE MATERIAL", size=11, bold=True, color=GOLD)
text(s, 8.75, 4.90, 3.8, 0.3, "AI ASSISTANCE", size=11, bold=True, color=GOLD)
text(s, 8.75, 5.24, 3.8, 1.6, [
    ("Claude \u2014 Sonnet and Opus (Anthropic)",
     {"size": 11.5, "bold": True, "space": 5}),
    ("Used throughout as a learning aid, coding mentor and analysis partner: study "
     "design review, Python analysis, figure generation, and iterative critique.",
     {"size": 10.5, "space": 5}),
    ("All measurements were run, and all conclusions checked, by the authors.",
     {"size": 10.5}),
], color=RGBColor(0xC8, 0xD2, 0xDC), line=1.35)
text(s, 0.85, 1.80, 7.5, 0.3, "LITERATURE", size=11, bold=True, color=GOLD)
notes(s, "PLACEHOLDER")


# =============================================================================
# 17 — contributions
# =============================================================================
s = new_slide(dark=True)
title(s, "Contributions", "who did what", dark=True)
for i, (initials, name, role) in enumerate([
        ("CK", "Christopher Kimberley",
         "derived sensor code, door-slam pilot protocol design, recording, "
         "and six test sessions"),
        ("CE", "Caitlin Everett",
         "product and study design, live tests, recording, export schema, analysis")]):
    yy = 3.05 + i * 1.60
    hexagon(s, 0.85, yy - 0.04, 0.68, initials, fill=GOLD, fg=NAVY, size=15)
    text(s, 1.90, yy, 10.5, 0.4, name, size=22, bold=True, color=WHITE)
    text(s, 1.90, yy + 0.52, 10.5, 0.5, role, size=16,
         color=RGBColor(0xC8, 0xD2, 0xDC))
notes(s, """
Who did what. The door-slam pilot — the protocol, the recording, and the six
sessions — is Christopher Kimberley's. The recorder, the export schema, the
analysis and the study design are mine.
""")

# =============================================================================
# 18 — close
# =============================================================================
s = new_slide(dark=True)
title(s, "Summary", dark=True)
REPO = "github.com/CaitlinEverett/ambient-recorder"
text(s, 0.85, 2.20, 8.9, 3.6, [
    ("The recorder works. All six channels run, and the derived vibration channel "
     "measures a real physical event at 13 to 109 times its own noise floor.",
     {"size": 20, "color": WHITE, "space": 15}),
    ("This is a feasibility result, not a reproducibility result. It is the best we "
     "could do in a tight window.",
     {"size": 20, "color": WHITE, "space": 15}),
    ("Whether a phone is a good enough instrument is still open. Whether it is a good "
     "enough recorder is not.",
     {"size": 20, "color": GOLD, "space": 15}),
    ("More studies run this week. The code and the data are public \u2014 scan the code, "
     "or:", {"size": 20, "color": WHITE, "space": 8}),
    (REPO, {"size": 20, "bold": True, "color": GOLD}),
], line=1.3)
_qr2 = FIG / "qr_repo_dark.png"
if _qr2.exists():
    s.shapes.add_picture(str(_qr2), Inches(10.05), Inches(4.55), Inches(2.15))
    text(s, 9.85, 6.82, 2.55, 0.3, "scan to contribute", size=12,
         color=RGBColor(0xC8, 0xD2, 0xDC), align=PP_ALIGN.CENTER)
notes(s, """
And the one that stung.

All six of those pilot sessions are labelled "controlled" — including both slams.
Our own app let us mislabel the entire dataset without a word of complaint.

We're building a tool to record what nobody wrote down, and it let us not write
something down. That's a real finding about the product, and it's in the report.

[CUE — cut back to the opening frame: the phone on the quiet counter. Hold two
seconds, then out.]
""")

# =============================================================================
# section dividers — authored last, positioned by ORDER
# =============================================================================
divider("1", "Aims and objectives",
        "what we set out to build, and the problem it is aimed at")
divider("2", "Project presentation",
        "what we built, and why consumer hardware makes it a question")
divider("3", "Changes to the plan",
        "four, and what each one cost or bought")
divider("4", "Results",
        "one pilot, three measured findings")
divider("5", "Reflection",
        "what was hard, what we got wrong, and what we do next")


# =============================================================================
# narration order — the single place slide sequence is decided
# =============================================================================
ORDER = [
    "__title__",
    "bluf1", "bluf2", "bluf3",
    "chris_app",
    "div:Aims and objectives",
    "idea0", "idea1", "idea2",
    "aims",
    "div:Project presentation",
    "Implementation",
    "inuse1", "inuse2", "inuse3", "inuse4",
    "Why this is hard",
    "Channels",
    "One test per sensor",
    "Cross-device alignment",
    "Privacy properties",
    "div:Changes to the plan",
    "Changes since the proposal",
    "div:Results",
    # The experiment beat is one compact block — the protocol card, then the
    # footage of it being run, then what came out. The n = 6 arithmetic moved
    # to Reflection: it is a lesson about designing a study, not a result.
    "The experiment, line by line",
    "Pilot study",
    "What the instrument measured",
    "Two phones, one table",
    "Derived vibration channel",
    "Effect of metric choice",
    "Sync fiducial recovery",
    "Unlabelled event detection",
    "div:Reflection",
    "Why n = 6",
    "What we got wrong",
    "Either outcome is useful",
    "Midway hypothesis",
    "What happens next",
    "Scope and limitations",
    "Standardising the disturbance",
    "Remaining work",
    "Works cited",
    "Contributions",
    "Summary",
]

REG["__title__"] = prs.slides[0]
missing = [k for k in ORDER if k not in REG]
extra = [k for k in REG if k not in ORDER]
if missing or extra:
    sys.exit(f"ORDER mismatch\n  missing from REG: {missing}\n  never ordered: {extra}")

_ids = {id(sl): e for sl, e in zip(prs.slides, list(sld_lst))}
for e in list(sld_lst):
    sld_lst.remove(e)
for k in ORDER:
    sld_lst.append(_ids[id(REG[k])])

# The lean deck: LEAN=1 removes the cut slides from the file instead of hiding
# them, so what is open on screen while filming is only what will be narrated.
# The full deck stays available for the written report.
import os

LEAN = os.environ.get("LEAN") == "1"
if LEAN:
    CUT = {
        "Channels", "One test per sensor", "Cross-device alignment",
        "Privacy properties", "Sync fiducial recovery", "Unlabelled event detection",
        "Scope and limitations", "Standardising the disturbance", "Remaining work",
        "Effect of metric choice", "Derived vibration channel",
        # Cut for time, not for weight. Eight minutes is a hard cap and the
        # spec mockup is superseded by the screenshot beside it; the
        # consumer-hardware failure-mode table survives in the full deck and
        # its conclusion — relative change only, never an absolute — is
        # already carried by the BLUF and by "Two phones, one table".
        "inuse1", "inuse2", "inuse3",
    }
    for k in CUT:
        e = _ids[id(REG[k])]
        prs.part.drop_rel(e.rId)
        sld_lst.remove(e)
    OUT = "Covariate_Demo.pptx"

prs.save(OUT)
print(f"wrote {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
